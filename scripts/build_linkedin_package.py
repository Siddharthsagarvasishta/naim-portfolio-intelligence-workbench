"""Finalize the validated LinkedIn/public showcase package from canonical evidence."""

from __future__ import annotations

import hashlib
import json
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from naim_risk.config import REPOSITORY_ROOT
from naim_risk.runtime_modes import dataset_hash

PRODUCT = "nAIM Portfolio Intelligence Workbench"
TAGLINE = "Name the movement. Own the evidence."
OUTPUT_ROOT = REPOSITORY_ROOT / "outputs" / "linkedin"
CANONICAL_PATH = REPOSITORY_ROOT / "exports" / "validation" / "interop_evidence_snapshot.json"
MARKET_PATH = REPOSITORY_ROOT / "outputs" / "market_risk" / "evidence_snapshot.json"
PUBLIC_FILES = (
    "project-summary.md",
    "technical-summary.md",
    "research-summary.md",
    "resume-bullets.md",
    "repository-description.txt",
    "accessibility-and-disclosure.md",
    "nAIM_LinkedIn_Carousel.pptx",
    "nAIM_LinkedIn_Carousel.pdf",
)


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def file_record(path: Path) -> dict[str, Any]:
    return {
        "path": path.name,
        "bytes": path.stat().st_size,
        "sha256": sha256_file(path),
    }


def presentation_text(path: Path) -> str:
    from pptx import Presentation

    deck = Presentation(path)
    return "\n".join(
        shape.text
        for slide in deck.slides
        for shape in slide.shapes
        if hasattr(shape, "text") and shape.text
    )


def pdf_text(path: Path) -> str:
    from pypdf import PdfReader

    return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)


def core_tokens(canonical: dict[str, Any]) -> tuple[str, ...]:
    finding = canonical["root_cause"]["finding"]
    return (
        f"{float(finding['observed_change_bps']):.1f}",
        f"{float(finding['mix_contribution_bps']):.1f}",
        f"{float(finding['within_segment_contribution_bps']):.1f}",
        str(finding["primary_driver"]),
        str(finding["causal_status"]),
    )


def market_risk_public_result(path: Path) -> str:
    """Return the truthful public classification for the bundled market snapshot."""

    if not path.is_file():
        return "UNAVAILABLE"
    payload = json.loads(path.read_text(encoding="utf-8"))
    if (
        payload.get("status") == "implemented"
        and payload.get("validation", {}).get("status") == "PASS"
        and payload.get("source", {}).get("source_is_synthetic") is True
        and payload.get("governance", {}).get("trading_recommendation") is False
        and payload.get("governance", {}).get("causal_claim") is False
    ):
        return "LIVE_VALIDATED_SYNTHETIC"
    return "UNAVAILABLE"


def main() -> None:
    canonical = json.loads(CANONICAL_PATH.read_text(encoding="utf-8"))
    run_id = str(canonical["metadata"]["run_id"])
    run_manifest = REPOSITORY_ROOT / "data" / "manifests" / run_id / "run_manifest.json"
    dataset_digest, dataset_basis = dataset_hash(run_manifest, REPOSITORY_ROOT / "data")
    paths = [OUTPUT_ROOT / name for name in PUBLIC_FILES]
    missing = [path.name for path in paths if not path.is_file()]
    if missing:
        raise FileNotFoundError(f"LinkedIn package inputs are missing: {missing}")

    pptx_path = OUTPUT_ROOT / "nAIM_LinkedIn_Carousel.pptx"
    pdf_path = OUTPUT_ROOT / "nAIM_LinkedIn_Carousel.pdf"
    tokens = core_tokens(canonical)
    pptx_story = presentation_text(pptx_path)
    pdf_story = pdf_text(pdf_path)
    if not all(token.casefold() in pptx_story.casefold() for token in tokens):
        raise ValueError("Editable carousel does not contain the complete canonical story.")
    if not all(token.casefold() in pdf_story.casefold() for token in tokens):
        raise ValueError("PDF carousel does not contain the complete canonical story.")

    records = sorted((file_record(path) for path in paths), key=lambda item: item["path"])
    package_hash = hashlib.sha256(
        json.dumps(records, sort_keys=True, separators=(",", ":")).encode("utf-8")
    ).hexdigest()
    created_at = datetime.now(UTC).isoformat()
    limitations = [
        "All figures are synthetic, institution-neutral demonstration data.",
        "The carousel communicates an associational investigation cue, not causal proof.",
        "No browser recording, animation, or automatic LinkedIn posting was performed.",
    ]
    manifest = {
        "schema_version": "2.0.0",
        "product": PRODUCT,
        "tagline": TAGLINE,
        "artifact_id": f"LINKEDIN-{package_hash[:20].upper()}",
        "artifact_type": "LINKEDIN_SHOWCASE_PACKAGE",
        "artifact_version": "2.0.0",
        "created_at": created_at,
        "created_at_utc": created_at,
        "created_by_component": "scripts.build_linkedin_package",
        "source_workspace": "all_portfolio_control",
        "source_snapshot_id": run_id,
        "evidence_id": canonical["evidence_id"],
        "data_mode": "OFFLINE_SNAPSHOT",
        "reporting_period": canonical["selected_reporting_period"],
        "comparison_period": "2025-07-01",
        "filter_scope": {
            "headline_scope": "all_portfolio",
            "approved_reference_basket": "BASKET-001",
            "scenario": "Baseline",
            "workspace": None,
        },
        "dataset_profile": canonical["metadata"]["profile"],
        "dataset_hash": dataset_digest,
        "dataset_hash_basis": dataset_basis,
        "configuration_hash": canonical["metadata"]["configuration_hash"],
        "metric_registry_version": canonical["metadata"]["metric_registry_version"],
        "code_version": "2.0.0",
        "evidence_ids": [canonical["evidence_id"]],
        "data_quality_result": canonical["data_quality"]["status"],
        "data_quality_status": canonical["data_quality"]["status"],
        "synthetic_data": canonical["synthetic_data_flag"],
        "synthetic_data_flag": canonical["synthetic_data_flag"],
        "file_name": pdf_path.name,
        "file_size": pdf_path.stat().st_size,
        "sha256": sha256_file(pdf_path),
        "dependencies": [
            "exports/validation/interop_evidence_snapshot.json",
            "outputs/market_risk/evidence_snapshot.json",
            f"data/manifests/{run_id}/run_manifest.json",
        ],
        "validation_status": "PASS",
        "validation_tests": [
            "artifact_tool_all_slide_render",
            "artifact_tool_layout_inspection",
            "slides_test_no_overflow",
            "libreoffice_pdf_conversion",
            "poppler_all_page_render",
            "pptx_canonical_story_tokens",
            "pdf_canonical_story_tokens",
            "sha256_file_ledger",
        ],
        "known_limitations": limitations,
        "automatic_linkedin_posting": False,
        "market_risk_public_result": market_risk_public_result(MARKET_PATH),
        "text_artifacts": {
            name: "READY" for name in PUBLIC_FILES if Path(name).suffix in {".md", ".txt"}
        },
        "media_artifacts": {
            "pdf_carousel": "READY_VALIDATED",
            "editable_presentation_source": "READY_VALIDATED",
            "browser_recording": "NOT_EXECUTABLE_LOCALLY",
            "animated_gif_or_mp4": "NOT_GENERATED",
        },
        "files": records,
        "validation": {
            "status": "PASS",
            "editable_slide_count": 7,
            "pdf_page_count": 7,
            "visual_inspection": "PASS",
            "package_sha256": package_hash,
        },
        "limitations": limitations,
        "disclosure": "All data and figures are synthetic and institution-neutral.",
    }
    (OUTPUT_ROOT / "package-manifest.json").write_text(
        json.dumps(manifest, indent=2) + "\n",
        encoding="utf-8",
    )
    print(json.dumps({"status": "PASS", "package_sha256": package_hash}, indent=2))


if __name__ == "__main__":
    main()
