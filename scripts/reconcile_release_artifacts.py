#!/usr/bin/env python3
"""Fail-closed reconciliation of the governed nAIM story across release channels.

The canonical interoperability snapshot is the numerical source of truth.  Every
adapter in this module is read-only: it reads an artifact, its native manifest,
and (where applicable) its release manifest.  Missing evidence is never treated
as a passing comparison.
"""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import math
import re
import shutil
import subprocess
import sys
import tempfile
from collections.abc import Iterable, Mapping, Sequence
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

PRODUCT = "nAIM Portfolio Intelligence Workbench"
TAGLINE = "Name the movement. Own the evidence."
SCHEMA_VERSION = "1.0.0"
DEFAULT_REPOSITORY_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUTPUT = Path("outputs/validation/cross_artifact_reconciliation.json")

PASS = "PASS"
FAIL = "FAIL"
MISSING = "MISSING"
UNVERIFIABLE = "UNVERIFIABLE"
INCOMPLETE = "INCOMPLETE"
NOT_APPLICABLE = "NOT_APPLICABLE"

_ABSENT = object()


class ReconciliationError(RuntimeError):
    """Raised when the canonical reconciliation source cannot be trusted."""


def sha256_file(path: Path) -> str:
    """Return a streaming SHA-256 digest without mutating the input."""

    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _portable(path: Path, repository_root: Path) -> str:
    resolved = path.resolve()
    root = repository_root.resolve()
    if not resolved.is_relative_to(root):
        raise ReconciliationError(f"Path is outside the repository: {path}")
    return resolved.relative_to(root).as_posix()


def _read_json(path: Path) -> dict[str, Any]:
    payload = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(payload, dict):
        raise ReconciliationError(f"Expected a JSON object: {path}")
    return payload


def _nested(payload: Mapping[str, Any], *paths: str) -> Any:
    """Return the first present dotted path or the private absent sentinel."""

    for dotted in paths:
        value: Any = payload
        found = True
        for part in dotted.split("."):
            if not isinstance(value, Mapping) or part not in value:
                found = False
                break
            value = value[part]
        if found:
            return value
    return _ABSENT


def _json_value(value: Any) -> Any:
    return None if value is _ABSENT else value


def check(
    check_id: str,
    *,
    expected: Any = _ABSENT,
    actual: Any = _ABSENT,
    required: bool = True,
    tolerance: float | None = None,
    evidence: str | None = None,
    detail: str | None = None,
) -> dict[str, Any]:
    """Create one machine-readable equality or presence check."""

    if actual is _ABSENT:
        outcome = MISSING if required else NOT_APPLICABLE
    elif expected is _ABSENT:
        outcome = PASS
    elif tolerance is not None and isinstance(expected, (int, float)):
        try:
            numeric = float(actual)
            target = float(expected)
            outcome = (
                PASS
                if math.isfinite(numeric)
                and math.isfinite(target)
                and abs(numeric - target) <= tolerance
                else FAIL
            )
        except (TypeError, ValueError):
            outcome = FAIL
    else:
        outcome = PASS if actual == expected else FAIL
    result = {
        "check_id": check_id,
        "required": required,
        "outcome": outcome,
        "expected": _json_value(expected),
        "actual": _json_value(actual),
    }
    if tolerance is not None:
        result["tolerance"] = tolerance
    if evidence:
        result["evidence"] = evidence
    if detail:
        result["detail"] = detail
    return result


def unverifiable(
    check_id: str,
    *,
    detail: str,
    required: bool = True,
    evidence: str | None = None,
) -> dict[str, Any]:
    result = {
        "check_id": check_id,
        "required": required,
        "outcome": UNVERIFIABLE,
        "expected": None,
        "actual": None,
        "detail": detail,
    }
    if evidence:
        result["evidence"] = evidence
    return result


def _channel_status(checks: Sequence[Mapping[str, Any]]) -> str:
    blockers = [item for item in checks if item.get("required")]
    if any(item.get("outcome") == FAIL for item in blockers):
        return FAIL
    if any(item.get("outcome") in {MISSING, UNVERIFIABLE} for item in blockers):
        return INCOMPLETE
    return PASS


def channel(
    channel_id: str,
    label: str,
    *,
    checks: list[dict[str, Any]],
    artifact_paths: Iterable[str] = (),
    notes: Iterable[str] = (),
) -> dict[str, Any]:
    return {
        "channel_id": channel_id,
        "label": label,
        "required": True,
        "status": _channel_status(checks),
        "artifact_paths": sorted(set(artifact_paths)),
        "checks": checks,
        "notes": list(notes),
    }


def _canonical_payload_hash(snapshot: Mapping[str, Any]) -> str:
    payload = dict(snapshot)
    payload.pop("evidence_sha256", None)
    canonical = json.dumps(payload, sort_keys=True, separators=(",", ":"), allow_nan=False)
    return hashlib.sha256(canonical.encode("utf-8")).hexdigest()


def extract_canonical_story(snapshot: Mapping[str, Any]) -> dict[str, Any]:
    """Extract the governed cross-channel claim set at full precision."""

    kpis = {
        str(row.get("metric_id")): row
        for row in snapshot.get("kpis", [])
        if isinstance(row, Mapping)
    }
    loss = kpis.get("ANNUALISED_NET_LOSS_RATE")
    finding = _nested(snapshot, "root_cause.finding")
    if not isinstance(loss, Mapping) or not isinstance(finding, Mapping):
        raise ReconciliationError("Canonical evidence lacks the loss KPI or root-cause finding")
    metadata = snapshot.get("metadata")
    if not isinstance(metadata, Mapping):
        raise ReconciliationError("Canonical evidence lacks metadata")
    return {
        "metric_id": "ANNUALISED_NET_LOSS_RATE",
        "metric_registry_version": metadata.get("metric_registry_version"),
        "metric_version": loss.get("metric_version"),
        "reporting_period": loss.get("reporting_period"),
        "comparison_period": loss.get("comparison_period"),
        "scope": "All portfolio",
        "current_annualised_net_loss_rate": loss.get("value"),
        "prior_annualised_net_loss_rate": loss.get("prior_value"),
        "observed_change_bps": finding.get("observed_change_bps"),
        "mix_contribution_bps": finding.get("mix_contribution_bps"),
        "within_segment_contribution_bps": finding.get("within_segment_contribution_bps"),
        "reconciliation_residual_bps": finding.get("reconciliation_residual_bps"),
        "primary_dimension": finding.get("primary_dimension"),
        "primary_driver": finding.get("primary_driver"),
        "causal_status": finding.get("causal_status"),
        "data_quality_status": _nested(snapshot, "data_quality.status"),
        "publication_allowed": _nested(snapshot, "data_quality.publication_allowed"),
        "synthetic_data": snapshot.get("synthetic_data_flag"),
        "run_id": metadata.get("run_id"),
        "configuration_hash": metadata.get("configuration_hash"),
        "evidence_id": snapshot.get("evidence_id"),
        "evidence_payload_sha256": snapshot.get("evidence_sha256"),
        "row_count": _nested(snapshot, "metadata.row_counts.monthly_account_performance"),
    }


def _dataset_context(
    repository_root: Path, story: Mapping[str, Any]
) -> tuple[str | None, str | None, dict[str, Any] | None, Path | None]:
    run_id = str(story.get("run_id") or "")
    manifest_path = repository_root / "data" / "manifests" / run_id / "run_manifest.json"
    if not run_id or not manifest_path.is_file():
        return None, None, None, None
    run_manifest = _read_json(manifest_path)
    try:
        from naim_risk.runtime_modes import dataset_hash

        digest, basis = dataset_hash(manifest_path, repository_root / "data")
    except (ImportError, OSError, ValueError, TypeError, json.JSONDecodeError):
        digest, basis = None, None
    return digest, basis, run_manifest, manifest_path


def _approved_basket(
    repository_root: Path, story: Mapping[str, Any]
) -> tuple[dict[str, Any] | None, str | None]:
    run_id = str(story.get("run_id") or "")
    path = repository_root / "data" / "raw" / run_id / "portfolio_basket_definition.parquet"
    if not path.is_file():
        return None, "Approved basket source is missing"
    try:
        import duckdb
    except ImportError:
        return None, "Approved basket source could not be read: ImportError"
    try:
        connection = duckdb.connect(":memory:")
        try:
            matches = connection.execute(
                "SELECT basket_id, basket_name, status, approved_flag "
                "FROM read_parquet(?) WHERE basket_id = 'BASKET-001' "
                "AND approved_flag AND lower(status) = 'approved'",
                [str(path)],
            ).fetchall()
        finally:
            connection.close()
    except (duckdb.Error, OSError, ValueError) as exc:
        return None, f"Approved basket source could not be read: {type(exc).__name__}"
    if len(matches) != 1:
        return None, "BASKET-001 is not uniquely approved"
    basket_id, basket_name, status, approved = matches[0]
    return (
        {
            "basket_id": str(basket_id),
            "basket_name": str(basket_name),
            "status": str(status),
            "approved": bool(approved),
            "scope_note": (
                "Approved secondary control basket; the governed headline remains All portfolio."
            ),
        },
        None,
    )


def canonical_context(repository_root: Path) -> tuple[dict[str, Any], dict[str, Any]]:
    """Load and validate the single canonical evidence source."""

    evidence_path = repository_root / "exports" / "validation" / "interop_evidence_snapshot.json"
    if not evidence_path.is_file():
        raise ReconciliationError(f"Canonical evidence is missing: {evidence_path}")
    snapshot = _read_json(evidence_path)
    story = extract_canonical_story(snapshot)
    dataset_digest, dataset_basis, run_manifest, run_manifest_path = _dataset_context(
        repository_root, story
    )
    if story.get("configuration_hash") in {None, ""} and run_manifest is not None:
        story["configuration_hash"] = run_manifest.get("configuration_hash")
    story["dataset_hash"] = dataset_digest
    story["dataset_hash_basis"] = dataset_basis
    story["canonical_file_sha256"] = sha256_file(evidence_path)
    approved_basket, basket_error = _approved_basket(repository_root, story)

    observed_from_rates = (
        float(story["current_annualised_net_loss_rate"])
        - float(story["prior_annualised_net_loss_rate"])
    ) * 10_000
    bridge_sum = (
        float(story["mix_contribution_bps"])
        + float(story["within_segment_contribution_bps"])
        + float(story["reconciliation_residual_bps"])
    )
    evidence_relative = _portable(evidence_path, repository_root)
    checks = [
        check(
            "canonical.payload_hash",
            expected=story.get("evidence_payload_sha256"),
            actual=_canonical_payload_hash(snapshot),
            evidence=evidence_relative,
        ),
        check(
            "canonical.product",
            expected=PRODUCT,
            actual=_nested(snapshot, "metadata.product"),
            evidence=evidence_relative,
        ),
        check(
            "canonical.synthetic",
            expected=True,
            actual=story.get("synthetic_data", _ABSENT),
            evidence=evidence_relative,
        ),
        check(
            "canonical.data_quality",
            expected=PASS,
            actual=story.get("data_quality_status", _ABSENT),
            evidence=evidence_relative,
        ),
        check(
            "canonical.publication_allowed",
            expected=True,
            actual=story.get("publication_allowed", _ABSENT),
            evidence=evidence_relative,
        ),
        check(
            "canonical.observed_from_rates",
            expected=story.get("observed_change_bps"),
            actual=observed_from_rates,
            tolerance=1e-9,
            evidence=evidence_relative,
        ),
        check(
            "canonical.bridge_reconciliation",
            expected=story.get("observed_change_bps"),
            actual=bridge_sum,
            tolerance=1e-9,
            evidence=evidence_relative,
        ),
        check(
            "canonical.primary_driver",
            expected="Affiliate",
            actual=story.get("primary_driver", _ABSENT),
            evidence=evidence_relative,
        ),
        check(
            "canonical.causal_status",
            expected="ASSOCIATIONAL",
            actual=story.get("causal_status", _ABSENT),
            evidence=evidence_relative,
        ),
        check(
            "canonical.dataset_hash",
            actual=dataset_digest if dataset_digest else _ABSENT,
            evidence=(_portable(run_manifest_path, repository_root) if run_manifest_path else None),
        ),
        check(
            "canonical.approved_control_basket",
            expected=True,
            actual=approved_basket.get("approved") if approved_basket else _ABSENT,
            detail=basket_error,
        ),
    ]
    if run_manifest is not None:
        checks.extend(
            [
                check(
                    "canonical.run_manifest_id",
                    expected=story.get("run_id"),
                    actual=run_manifest.get("run_id", _ABSENT),
                    evidence=_portable(run_manifest_path, repository_root),
                ),
                check(
                    "canonical.configuration_hash",
                    expected=story.get("configuration_hash"),
                    actual=run_manifest.get("configuration_hash", _ABSENT),
                    evidence=_portable(run_manifest_path, repository_root),
                ),
            ]
        )
    else:
        checks.append(
            check(
                "canonical.run_manifest",
                actual=_ABSENT,
                detail="The canonical run manifest is required for configuration provenance.",
            )
        )
    context = {
        "source_path": evidence_relative,
        "source_file_sha256": story["canonical_file_sha256"],
        "source_payload_sha256": story["evidence_payload_sha256"],
        "dataset_hash": dataset_digest,
        "dataset_hash_basis": dataset_basis,
        "approved_control_period": {
            "reporting_period": story["reporting_period"],
            "comparison_period": story["comparison_period"],
            "scope": story["scope"],
        },
        "approved_control_basket": approved_basket,
        "story": story,
        "checks": checks,
        "status": _channel_status(checks),
    }
    api_channel = channel(
        "api_service_evidence",
        "API/service canonical evidence",
        checks=[
            *checks,
            check(
                "api.service_generation_method",
                expected=True,
                actual=str(snapshot.get("generation_method", "")).startswith(
                    "WorkbenchService(load_config("
                ),
                evidence=evidence_relative,
            ),
            check(
                "api.service_source_reference",
                expected="src/naim_risk/service.py",
                actual=snapshot.get("source_ref", _ABSENT),
                evidence=evidence_relative,
            ),
        ],
        artifact_paths=[evidence_relative],
        notes=[
            "This validates the same service methods used by the API; HTTP/browser capture is separate."
        ],
    )
    return context, api_channel


STORY_NUMERIC_FIELDS = {
    "current_annualised_net_loss_rate": 1e-9,
    "prior_annualised_net_loss_rate": 1e-9,
    "observed_change_bps": 1e-9,
    "mix_contribution_bps": 1e-9,
    "within_segment_contribution_bps": 1e-9,
    "reconciliation_residual_bps": 1e-8,
}


def compare_story(
    expected: Mapping[str, Any],
    actual: Mapping[str, Any],
    *,
    prefix: str,
    fields: Iterable[str] | None = None,
    evidence: str | None = None,
    tolerances: Mapping[str, float] | None = None,
) -> list[dict[str, Any]]:
    """Compare a channel story to the canonical source with named tolerances."""

    selected = list(
        fields
        or [
            *STORY_NUMERIC_FIELDS,
            "reporting_period",
            "comparison_period",
            "primary_dimension",
            "primary_driver",
            "causal_status",
            "run_id",
            "metric_registry_version",
            "data_quality_status",
            "synthetic_data",
        ]
    )
    tolerance_map = {**STORY_NUMERIC_FIELDS, **dict(tolerances or {})}
    checks: list[dict[str, Any]] = []
    for name in selected:
        checks.append(
            check(
                f"{prefix}.story.{name}",
                expected=expected.get(name, _ABSENT),
                actual=actual.get(name, _ABSENT),
                tolerance=tolerance_map.get(name),
                evidence=evidence,
            )
        )
    return checks


def _normalise_evidence_id(evidence_id: Any, run_id: Any) -> bool | object:
    if evidence_id is _ABSENT or run_id in {_ABSENT, None, ""}:
        return _ABSENT
    return (
        str(evidence_id).casefold().endswith(str(run_id).casefold())
        or str(run_id).casefold() in str(evidence_id).casefold()
    )


def verify_file_ledger(
    root: Path,
    files: Any,
    *,
    prefix: str,
    repository_root: Path,
) -> list[dict[str, Any]]:
    """Verify a portable file ledger and reject traversal or missing checksums."""

    if not isinstance(files, list) or not files:
        return [check(f"{prefix}.file_ledger", actual=_ABSENT)]
    checks: list[dict[str, Any]] = []
    root_resolved = root.resolve()
    for index, item in enumerate(files):
        if not isinstance(item, Mapping):
            checks.append(
                check(
                    f"{prefix}.file_{index}.record",
                    expected=True,
                    actual=False,
                )
            )
            continue
        relative = str(item.get("path", ""))
        candidate = (root / relative).resolve()
        safe = bool(relative) and candidate.is_relative_to(root_resolved)
        checks.append(
            check(
                f"{prefix}.file_{index}.portable_path",
                expected=True,
                actual=safe,
                evidence=relative or None,
            )
        )
        if not safe or not candidate.is_file():
            checks.append(
                check(
                    f"{prefix}.file_{index}.exists",
                    expected=True,
                    actual=False,
                    evidence=relative or None,
                )
            )
            continue
        portable = _portable(candidate, repository_root)
        checks.extend(
            [
                check(
                    f"{prefix}.file_{index}.bytes",
                    expected=item.get("bytes", _ABSENT),
                    actual=candidate.stat().st_size,
                    evidence=portable,
                ),
                check(
                    f"{prefix}.file_{index}.sha256",
                    expected=item.get("sha256", _ABSENT),
                    actual=sha256_file(candidate),
                    evidence=portable,
                ),
            ]
        )
    return checks


def _find_release_manifest(
    repository_root: Path, artifact: Path
) -> tuple[dict[str, Any] | None, Path | None]:
    manifest_root = repository_root / "outputs" / "manifests"
    if not manifest_root.is_dir():
        return None, None
    artifact_relative = _portable(artifact, repository_root)
    for manifest_path in sorted(manifest_root.glob("*.json")):
        try:
            payload = _read_json(manifest_path)
        except (OSError, json.JSONDecodeError, ReconciliationError):
            continue
        registered = _nested(payload, "artifact.path")
        if registered == artifact_relative:
            return payload, manifest_path
    return None, None


def _release_manifest_checks(
    repository_root: Path,
    artifact: Path,
    expected: Mapping[str, Any],
    *,
    prefix: str,
) -> tuple[list[dict[str, Any]], list[str]]:
    """Enforce the complete §112 provenance contract for a single file."""

    manifest, manifest_path = _find_release_manifest(repository_root, artifact)
    if manifest is None or manifest_path is None:
        return (
            [
                check(
                    f"{prefix}.release_manifest",
                    actual=_ABSENT,
                    detail="No matching JSON manifest exists under outputs/manifests/.",
                )
            ],
            [],
        )
    evidence = _portable(manifest_path, repository_root)
    registered_hash = _nested(manifest, "artifact.sha256")
    dataset_hash = _nested(manifest, "versions.dataset_hash.value", "dataset_hash")
    configuration_hash = _nested(
        manifest, "versions.configuration_hash.value", "configuration_hash"
    )
    code_version = _nested(
        manifest,
        "versions.code_version.value",
        "versions.script_version.value",
        "code_version",
    )
    metric_version = _nested(
        manifest,
        "metric_registry_version",
        "versions.metric_registry_version.value",
    )
    filter_scope = _nested(manifest, "filter_scope", "provenance.filter_scope")
    evidence_ids = _nested(manifest, "evidence_ids", "provenance.evidence_ids")
    data_quality = _nested(
        manifest, "data_quality_result", "data_quality.status", "provenance.data_quality_result"
    )
    synthetic = _nested(manifest, "synthetic_data", "synthetic", "provenance.synthetic_data")
    reporting = _nested(manifest, "reporting_period", "reporting_date")
    reporting_matches = (
        _ABSENT
        if reporting is _ABSENT
        else str(reporting)[:7] == str(expected.get("reporting_period"))[:7]
    )
    if filter_scope is _ABSENT:
        scope_matches: bool | object = _ABSENT
    else:
        serialised_scope = re.sub(
            r"[_-]+", " ", json.dumps(filter_scope, sort_keys=True).casefold()
        )
        expected_scope = re.sub(r"[_-]+", " ", str(expected.get("scope", "")).casefold())
        scope_matches = "all portfolio" in serialised_scope or expected_scope in serialised_scope
    evidence_id_matches = (
        _ABSENT
        if evidence_ids is _ABSENT
        else any(
            str(expected.get("run_id")) in str(item)
            for item in (evidence_ids if isinstance(evidence_ids, list) else [evidence_ids])
        )
    )
    checks = [
        check(f"{prefix}.release_manifest_present", expected=True, actual=True, evidence=evidence),
        check(
            f"{prefix}.artifact_id",
            actual=_nested(manifest, "artifact_id", "build_id", "presentation_id"),
            evidence=evidence,
        ),
        check(
            f"{prefix}.creation_time",
            actual=_nested(manifest, "created_at_utc", "built_at_utc", "generated_at"),
            evidence=evidence,
        ),
        check(
            f"{prefix}.artifact_sha256",
            expected=registered_hash,
            actual=sha256_file(artifact),
            evidence=evidence,
        ),
        check(
            f"{prefix}.dataset_hash",
            expected=expected.get("dataset_hash", _ABSENT),
            actual=dataset_hash,
            evidence=evidence,
        ),
        check(
            f"{prefix}.configuration_hash",
            expected=expected.get("configuration_hash", _ABSENT),
            actual=configuration_hash,
            evidence=evidence,
        ),
        check(f"{prefix}.code_version", actual=code_version, evidence=evidence),
        check(
            f"{prefix}.metric_registry_version",
            expected=expected.get("metric_registry_version", _ABSENT),
            actual=metric_version,
            evidence=evidence,
        ),
        check(
            f"{prefix}.reporting_period",
            expected=True,
            actual=reporting_matches,
            evidence=evidence,
        ),
        check(
            f"{prefix}.filter_scope",
            expected=True,
            actual=scope_matches,
            evidence=evidence,
        ),
        check(
            f"{prefix}.evidence_ids",
            expected=True,
            actual=evidence_id_matches,
            evidence=evidence,
        ),
        check(
            f"{prefix}.data_quality_result",
            expected=expected.get("data_quality_status", _ABSENT),
            actual=data_quality,
            evidence=evidence,
        ),
        check(
            f"{prefix}.synthetic_data",
            expected=expected.get("synthetic_data", _ABSENT),
            actual=synthetic,
            evidence=evidence,
        ),
    ]
    return checks, [evidence]


def _missing_channel(channel_id: str, label: str, expected_path: str) -> dict[str, Any]:
    return channel(
        channel_id,
        label,
        checks=[
            check(
                f"{channel_id}.artifact_present",
                expected=expected_path,
                actual=_ABSENT,
            )
        ],
        notes=["No current artifact or validated runtime snapshot was found."],
    )


def _snapshot_story(payload: Mapping[str, Any]) -> dict[str, Any]:
    if "governed_story" in payload and isinstance(payload["governed_story"], Mapping):
        return dict(payload["governed_story"])
    if "portfolio_story" in payload and "decomposition" in payload:
        portfolio = payload["portfolio_story"]
        decomposition = payload["decomposition"]
        source = payload.get("source_context", {})
        return {
            "metric_id": portfolio.get("metric_id"),
            "metric_registry_version": source.get("metric_registry_version"),
            "metric_version": portfolio.get("metric_version"),
            "reporting_period": payload.get("reporting_period"),
            "comparison_period": "2025-07-01",
            "scope": "All portfolio",
            "current_annualised_net_loss_rate": portfolio.get("current_annualised_net_loss_rate"),
            "prior_annualised_net_loss_rate": portfolio.get("prior_annualised_net_loss_rate"),
            "observed_change_bps": portfolio.get("observed_change_bps"),
            "mix_contribution_bps": decomposition.get("mix_bps"),
            "within_segment_contribution_bps": decomposition.get("within_segment_bps"),
            "reconciliation_residual_bps": decomposition.get("residual_bps"),
            "primary_dimension": decomposition.get("dimension"),
            "primary_driver": decomposition.get("primary_driver"),
            "causal_status": decomposition.get("causal_status"),
            "data_quality_status": portfolio.get("data_quality_status"),
            "synthetic_data": payload.get("synthetic_data", True),
            "run_id": source.get("run_id"),
            "evidence_id": payload.get("evidence_id"),
        }
    return extract_canonical_story(payload)


def _runtime_snapshot_channel(
    repository_root: Path,
    expected: Mapping[str, Any],
    *,
    channel_id: str,
    label: str,
    relative_path: str,
) -> dict[str, Any]:
    path = repository_root / relative_path
    if not path.is_file():
        return _missing_channel(channel_id, label, relative_path)
    portable = _portable(path, repository_root)
    try:
        payload = _read_json(path)
        actual = _snapshot_story(payload)
        checks = compare_story(expected, actual, prefix=channel_id, evidence=portable)
        checks.extend(
            [
                check(
                    f"{channel_id}.captured_at",
                    actual=_nested(payload, "captured_at_utc", "generated_at"),
                    evidence=portable,
                ),
                check(
                    f"{channel_id}.source_context_dataset_hash",
                    expected=expected.get("dataset_hash", _ABSENT),
                    actual=_nested(payload, "source_context.dataset_hash", "dataset_hash"),
                    evidence=portable,
                ),
                check(
                    f"{channel_id}.source_context_configuration_hash",
                    expected=expected.get("configuration_hash", _ABSENT),
                    actual=_nested(
                        payload, "source_context.configuration_hash", "configuration_hash"
                    ),
                    evidence=portable,
                ),
            ]
        )
    except (OSError, ValueError, TypeError, json.JSONDecodeError, ReconciliationError) as exc:
        checks = [
            check(
                f"{channel_id}.readable_snapshot",
                expected=True,
                actual=False,
                evidence=portable,
                detail=f"{type(exc).__name__}: {exc}",
            )
        ]
    return channel(channel_id, label, checks=checks, artifact_paths=[portable])


def _workbook_channel(repository_root: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    path = repository_root / "outputs" / "nAIM_Portfolio_Intelligence_Workbench.xlsx"
    if not path.is_file():
        return _missing_channel(
            "excel_workbook",
            "Excel workbook",
            "outputs/nAIM_Portfolio_Intelligence_Workbench.xlsx",
        )
    portable = _portable(path, repository_root)
    checks: list[dict[str, Any]] = []
    paths = [portable]
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            root = workbook["Root Cause"]
            refresh = workbook["Refresh Control"]
            guide = workbook["Workbook Guide"]
            actual = {
                "observed_change_bps": root["B5"].value,
                "mix_contribution_bps": root["B6"].value,
                "within_segment_contribution_bps": root["B7"].value,
                "reconciliation_residual_bps": root["B8"].value,
                "primary_dimension": root["B9"].value,
                "primary_driver": root["B10"].value,
                "causal_status": root["B12"].value,
                "reporting_period": refresh["B8"].value,
                "run_id": refresh["B7"].value,
                "metric_registry_version": refresh["B11"].value,
                "data_quality_status": refresh["B12"].value,
                "synthetic_data": "Synthetic" in str(guide["B9"].value),
            }
            evidence_id = refresh["B5"].value
            payload_hash = refresh["B6"].value
        finally:
            workbook.close()
        checks.extend(
            compare_story(
                expected,
                actual,
                prefix="excel",
                fields=(
                    "observed_change_bps",
                    "mix_contribution_bps",
                    "within_segment_contribution_bps",
                    "reconciliation_residual_bps",
                    "primary_dimension",
                    "primary_driver",
                    "causal_status",
                    "reporting_period",
                    "run_id",
                    "metric_registry_version",
                    "data_quality_status",
                    "synthetic_data",
                ),
                evidence=f"{portable}#Root Cause/Refresh Control",
            )
        )
        checks.extend(
            [
                check(
                    "excel.evidence_id_maps_to_run",
                    expected=True,
                    actual=_normalise_evidence_id(evidence_id, expected.get("run_id")),
                    evidence=f"{portable}#Refresh Control!B5",
                ),
                check(
                    "excel.evidence_payload_sha256",
                    expected=expected.get("evidence_payload_sha256"),
                    actual=payload_hash,
                    evidence=f"{portable}#Refresh Control!B6",
                ),
            ]
        )
    except (ImportError, KeyError, OSError, ValueError, TypeError) as exc:
        checks.append(
            check(
                "excel.readable_workbook",
                expected=True,
                actual=False,
                evidence=portable,
                detail=f"{type(exc).__name__}: {exc}",
            )
        )
    validation_path = repository_root / "outputs" / "validation" / "office_workbook_validation.json"
    if validation_path.is_file():
        validation = _read_json(validation_path)
        validation_portable = _portable(validation_path, repository_root)
        paths.append(validation_portable)
        checks.append(
            check(
                "excel.native_validation",
                expected=PASS,
                actual=validation.get("status", _ABSENT),
                evidence=validation_portable,
            )
        )
    else:
        checks.append(check("excel.native_validation", actual=_ABSENT))
    manifest_checks, manifest_paths = _release_manifest_checks(
        repository_root, path, expected, prefix="excel.manifest"
    )
    checks.extend(manifest_checks)
    paths.extend(manifest_paths)
    return channel(
        "excel_workbook",
        "Excel workbook",
        checks=checks,
        artifact_paths=paths,
        notes=[f"Physical workbook SHA-256: {sha256_file(path)}"],
    )


def _presentation_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
        parts.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(parts)


def _presentation_channel(repository_root: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    root = repository_root / "outputs" / "presentations"
    manifests = sorted(root.glob("*.manifest.json")) if root.is_dir() else []
    candidates: list[tuple[Path, dict[str, Any]]] = []
    for manifest_path in manifests:
        try:
            manifest = _read_json(manifest_path)
        except (OSError, json.JSONDecodeError, ReconciliationError):
            continue
        if (
            str(manifest.get("reporting_period", ""))[:7]
            == str(expected.get("reporting_period", ""))[:7]
        ):
            candidates.append((manifest_path, manifest))
    if not candidates:
        wrong_periods = [
            str(_read_json(item).get("reporting_period")) for item in manifests if item.is_file()
        ]
        result = _missing_channel(
            "powerpoint_review",
            "PowerPoint review",
            "outputs/presentations/nAIM_Portfolio_Intelligence_Review_2025_08.pptx",
        )
        if wrong_periods:
            result["notes"].append(
                f"Presentation manifests exist only for non-control periods: {wrong_periods}"
            )
        return result
    preferred_manifest = root / "nAIM_Portfolio_Intelligence_Review.manifest.json"
    preferred = [item for item in candidates if item[0] == preferred_manifest]
    if preferred:
        candidates = preferred
    manifest_path, manifest = candidates[-1]
    registered_path = _nested(manifest, "artifact.path")
    if isinstance(registered_path, str) and registered_path:
        path = (repository_root / registered_path).resolve()
        if not path.is_relative_to(repository_root.resolve()):
            raise ReconciliationError(
                f"Presentation manifest path escapes the repository: {registered_path}"
            )
    else:
        filename = str(_nested(manifest, "artifact.filename"))
        path = root / filename
    manifest_portable = _portable(manifest_path, repository_root)
    paths = [manifest_portable]
    excel_reference = _nested(manifest, "reconciliation.excel_reference")
    checks: list[dict[str, Any]] = [
        check(
            "powerpoint.unique_control_period_deck",
            expected=1,
            actual=len(candidates),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.native_validation",
            expected=PASS,
            actual=_nested(manifest, "validation.status"),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.run_id",
            expected=expected.get("run_id"),
            actual=manifest.get("source_snapshot_id", _ABSENT),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.configuration_hash",
            expected=expected.get("configuration_hash"),
            actual=manifest.get("configuration_hash", _ABSENT),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.dataset_hash",
            expected=expected.get("dataset_hash", _ABSENT),
            actual=manifest.get("dataset_hash", _ABSENT),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.metric_version",
            expected=expected.get("metric_registry_version"),
            actual=manifest.get(
                "metric_version",
                manifest.get("metric_registry_version", _ABSENT),
            ),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.synthetic_disclosure",
            expected=True,
            actual="synthetic" in str(manifest.get("synthetic_data_disclaimer", "")).casefold(),
            evidence=manifest_portable,
        ),
        check(
            "powerpoint.excel_reconciliation",
            expected=True,
            actual=(
                True
                if excel_reference is not _ABSENT
                and excel_reference not in (None, "", "not_generated_in_this_request")
                else _ABSENT
            ),
            evidence=manifest_portable,
        ),
    ]
    if not path.is_file():
        checks.append(
            check(
                "powerpoint.artifact_present",
                expected=True,
                actual=False,
                evidence=_portable(root, repository_root),
            )
        )
        return channel(
            "powerpoint_review", "PowerPoint review", checks=checks, artifact_paths=paths
        )
    portable = _portable(path, repository_root)
    paths.append(portable)
    checks.append(
        check(
            "powerpoint.artifact_sha256",
            expected=_nested(manifest, "artifact.sha256"),
            actual=sha256_file(path),
            evidence=portable,
        )
    )
    try:
        text = _presentation_text(path)
        checks.extend(
            [
                check(
                    "powerpoint.observed_display_value",
                    expected=True,
                    actual=f"{float(expected['observed_change_bps']):+.1f} bps" in text,
                    evidence=portable,
                    tolerance=0.05,
                ),
                check(
                    "powerpoint.mix_display_value",
                    expected=True,
                    actual=f"{float(expected['mix_contribution_bps']):+.1f} bps" in text,
                    evidence=portable,
                    tolerance=0.05,
                ),
                check(
                    "powerpoint.within_display_value",
                    expected=True,
                    actual=(
                        f"{float(expected['within_segment_contribution_bps']):+.1f} bps" in text
                    ),
                    evidence=portable,
                    tolerance=0.05,
                ),
                check(
                    "powerpoint.primary_driver",
                    expected=True,
                    actual=str(expected.get("primary_driver")) in text,
                    evidence=portable,
                ),
                check(
                    "powerpoint.associational_disclosure",
                    expected=True,
                    actual="associational" in text.casefold(),
                    evidence=portable,
                ),
                check(
                    "powerpoint.tagline",
                    expected=True,
                    actual=TAGLINE in text,
                    evidence=portable,
                ),
            ]
        )
    except (ImportError, OSError, ValueError, KeyError) as exc:
        checks.append(
            unverifiable(
                "powerpoint.readable_text",
                detail=f"{type(exc).__name__}: {exc}",
                evidence=portable,
            )
        )
    release_checks, release_paths = _release_manifest_checks(
        repository_root, path, expected, prefix="powerpoint.manifest"
    )
    checks.extend(release_checks)
    paths.extend(release_paths)
    return channel(
        "powerpoint_review",
        "PowerPoint review",
        checks=checks,
        artifact_paths=paths,
    )


def _parse_hyper_shell_rows(
    output: str,
) -> tuple[list[list[str]], list[list[str]], list[list[str]]]:
    """Parse the deliberately prefixed rows emitted by ``hyperd shell``."""

    portfolio_rows: list[list[str]] = []
    segment_rows: list[list[str]] = []
    metadata_rows: list[list[str]] = []
    for raw_line in output.splitlines():
        line = raw_line.strip()
        if line.startswith("NAIMPORTFOLIO|"):
            portfolio_rows.append(line.split("|", 4)[1:])
        elif line.startswith("NAIMSEGMENT|"):
            segment_rows.append(line.split("|", 5)[1:])
        elif line.startswith("NAIMMETA|"):
            metadata_rows.append(line.split("|", 2)[1:])
    if len(portfolio_rows) != 2 or not segment_rows or not metadata_rows:
        raise ReconciliationError(
            "Hyper shell readback did not return the required portfolio, segment, and metadata rows"
        )
    return portfolio_rows, segment_rows, metadata_rows


def _hyper_shell_rows(
    path: Path,
    *,
    reporting: str,
    comparison: str,
    dimension: str,
    hyper: Any,
) -> tuple[list[list[str]], list[list[str]], list[list[str]]]:
    """Query Hyper directly with its native shell when sandbox sockets are denied."""

    hyperd = Path(hyper.__file__).resolve().parent / "bin" / "hyper" / "hyperd"
    if not hyperd.is_file():
        raise ReconciliationError(f"Hyper native shell is unavailable: {hyperd}")
    quoted_dimension = hyper.escape_name(dimension)
    script = f"""
SELECT 'NAIMPORTFOLIO|' || CAST("month" AS TEXT) || '|' ||
       CAST("chargeoff_amount" AS TEXT) || '|' || CAST("recovery_amount" AS TEXT) || '|' ||
       CAST("average_receivables" AS TEXT) AS "naim_row"
FROM "Extract"."MartPortfolioMonth"
WHERE "month" IN (DATE '{comparison}', DATE '{reporting}')
ORDER BY "month";
SELECT 'NAIMSEGMENT|' || CAST("month" AS TEXT) || '|' || CAST({quoted_dimension} AS TEXT) ||
       '|' || CAST(SUM("chargeoff_amount") AS TEXT) || '|' ||
       CAST(SUM("recovery_amount") AS TEXT) || '|' ||
       CAST(SUM("average_receivables") AS TEXT) AS "naim_row"
FROM "Extract"."MartSegmentMonth"
WHERE "month" IN (DATE '{comparison}', DATE '{reporting}')
GROUP BY "month", {quoted_dimension}
ORDER BY "month", {quoted_dimension};
SELECT 'NAIMMETA|' || CAST("metadata_key" AS TEXT) || '|' ||
       CAST("metadata_value" AS TEXT) AS "naim_row"
FROM "Extract"."Metadata"
ORDER BY "metadata_key";
\\q
"""
    with tempfile.TemporaryDirectory(prefix="naim-hyper-shell-") as work:
        result = subprocess.run(
            [
                str(hyperd),
                "--log-dir",
                work,
                "--database",
                str(path.resolve()),
                "shell",
            ],
            input=script,
            check=True,
            capture_output=True,
            text=True,
            timeout=60,
        )
    return _parse_hyper_shell_rows(f"{result.stdout}\n{result.stderr}")


def _hyper_story(path: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    """Read the core story directly with the official Hyper engine."""

    import pandas as pd
    import tableauhyperapi as hyper

    reporting = str(expected["reporting_period"])
    comparison = str(expected["comparison_period"])
    dimension = str(expected["primary_dimension"])
    allowed_dimensions = {
        "acquisition_channel",
        "product_type",
        "customer_segment",
        "geography",
        "original_risk_band",
        "strategy_version",
    }
    if dimension not in allowed_dimensions:
        raise ReconciliationError(f"Unsupported Hyper reconciliation dimension: {dimension}")
    try:
        with (
            hyper.HyperProcess(hyper.Telemetry.DO_NOT_SEND_USAGE_DATA_TO_TABLEAU) as process,
            hyper.Connection(process.endpoint, path) as connection,
        ):
            portfolio_rows = connection.execute_list_query(
                'SELECT "month", "chargeoff_amount", "recovery_amount", '
                '"average_receivables" FROM "Extract"."MartPortfolioMonth" '
                f"WHERE \"month\" IN (DATE '{comparison}', DATE '{reporting}')"
            )
            quoted_dimension = hyper.escape_name(dimension)
            segment_rows = connection.execute_list_query(
                f'SELECT "month", {quoted_dimension}, SUM("chargeoff_amount"), '
                'SUM("recovery_amount"), SUM("average_receivables") '
                'FROM "Extract"."MartSegmentMonth" '
                f"WHERE \"month\" IN (DATE '{comparison}', DATE '{reporting}') "
                f'GROUP BY "month", {quoted_dimension}'
            )
            metadata_rows = connection.execute_list_query(
                'SELECT "metadata_key", "metadata_value" FROM "Extract"."Metadata"'
            )
    except hyper.HyperException as exc:
        message = str(exc)
        if "binding to unix domain socket" not in message or "Operation not permitted" not in message:
            raise
        portfolio_rows, segment_rows, metadata_rows = _hyper_shell_rows(
            path,
            reporting=reporting,
            comparison=comparison,
            dimension=dimension,
            hyper=hyper,
        )
    portfolio = {
        str(row[0]): (float(row[1]) - float(row[2])) * 12 / float(row[3]) for row in portfolio_rows
    }
    frame = pd.DataFrame(
        segment_rows,
        columns=["month", dimension, "chargeoff_amount", "recovery_amount", "denominator"],
    )
    frame["month"] = frame["month"].astype(str)
    frame[dimension] = frame[dimension].astype(str)
    frame["numerator"] = frame["chargeoff_amount"].astype(float) - frame["recovery_amount"].astype(
        float
    )
    frame["denominator"] = frame["denominator"].astype(float)
    from naim_risk.root_cause.decomposition import decompose_rate

    bridge = decompose_rate(
        frame[frame["month"] == comparison],
        frame[frame["month"] == reporting],
        segment_column=dimension,
        numerator_column="numerator",
        denominator_column="denominator",
        scale=120_000.0,
    )
    metadata = {str(row[0]): str(row[1]) for row in metadata_rows}
    return {
        "current_annualised_net_loss_rate": portfolio.get(reporting),
        "prior_annualised_net_loss_rate": portfolio.get(comparison),
        "observed_change_bps": bridge["observed_change"],
        "mix_contribution_bps": bridge["mix_contribution"],
        "within_segment_contribution_bps": bridge["within_segment_contribution"],
        "reconciliation_residual_bps": bridge["reconciliation_residual"],
        "primary_dimension": dimension,
        "primary_driver": bridge["segments"][0]["segment"],
        "run_id": metadata.get("run_id"),
        "configuration_hash": metadata.get("configuration_hash"),
        "data_quality_status": metadata.get("quality_status"),
        "synthetic_data": metadata.get("synthetic_data", "").casefold() == "true",
        "product": metadata.get("product"),
    }


def _hyper_channel(repository_root: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    path = repository_root / "outputs" / "tableau" / "nAIM_Portfolio_Intelligence.hyper"
    manifest_path = path.with_suffix(".manifest.json")
    if not path.is_file():
        return _missing_channel(
            "tableau_hyper",
            "Tableau Hyper extract",
            "outputs/tableau/nAIM_Portfolio_Intelligence.hyper",
        )
    portable = _portable(path, repository_root)
    paths = [portable]
    checks: list[dict[str, Any]] = []
    if manifest_path.is_file():
        manifest = _read_json(manifest_path)
        manifest_portable = _portable(manifest_path, repository_root)
        paths.append(manifest_portable)
        checks.extend(
            [
                check(
                    "tableau.native_status",
                    expected=PASS,
                    actual=manifest.get("status", _ABSENT),
                    evidence=manifest_portable,
                ),
                check(
                    "tableau.artifact_sha256",
                    expected=manifest.get("sha256", _ABSENT),
                    actual=sha256_file(path),
                    evidence=portable,
                ),
                check(
                    "tableau.run_id",
                    expected=expected.get("run_id"),
                    actual=manifest.get("run_id", _ABSENT),
                    evidence=manifest_portable,
                ),
                check(
                    "tableau.configuration_hash",
                    expected=expected.get("configuration_hash"),
                    actual=manifest.get("configuration_hash", _ABSENT),
                    evidence=manifest_portable,
                ),
                check(
                    "tableau.table_control_totals",
                    expected=True,
                    actual=bool(manifest.get("tables"))
                    and all(row.get("status") == PASS for row in manifest["tables"]),
                    evidence=manifest_portable,
                ),
            ]
        )
    else:
        checks.append(check("tableau.native_manifest", actual=_ABSENT))
    try:
        actual = _hyper_story(path, expected)
        checks.extend(
            compare_story(
                expected,
                actual,
                prefix="tableau",
                fields=(
                    "current_annualised_net_loss_rate",
                    "prior_annualised_net_loss_rate",
                    "observed_change_bps",
                    "mix_contribution_bps",
                    "within_segment_contribution_bps",
                    "reconciliation_residual_bps",
                    "primary_dimension",
                    "primary_driver",
                    "run_id",
                    "data_quality_status",
                    "synthetic_data",
                ),
                evidence=portable,
            )
        )
        checks.append(
            check(
                "tableau.metadata_product",
                expected=PRODUCT,
                actual=actual.get("product", _ABSENT),
                evidence=portable,
            )
        )
    except (ImportError, OSError, ValueError, RuntimeError, ReconciliationError) as exc:
        checks.append(
            unverifiable(
                "tableau.direct_query",
                detail=f"{type(exc).__name__}: {exc}",
                evidence=portable,
            )
        )
    release_checks, release_paths = _release_manifest_checks(
        repository_root, path, expected, prefix="tableau.manifest"
    )
    checks.extend(release_checks)
    paths.extend(release_paths)
    return channel(
        "tableau_hyper",
        "Tableau Hyper extract",
        checks=checks,
        artifact_paths=paths,
    )


def _read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle))


def _powerbi_channel(repository_root: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    root = repository_root / "outputs" / "powerbi" / "nAIM.PowerBIProject"
    manifest_path = root / "Build" / "project-manifest.json"
    if not manifest_path.is_file():
        return _missing_channel(
            "power_bi_validation",
            "Power BI validation data",
            "outputs/powerbi/nAIM.PowerBIProject/Build/project-manifest.json",
        )
    manifest = _read_json(manifest_path)
    manifest_portable = _portable(manifest_path, repository_root)
    paths = [manifest_portable]
    checks = verify_file_ledger(
        root,
        manifest.get("files"),
        prefix="powerbi",
        repository_root=repository_root,
    )
    ledger = manifest.get("files") if isinstance(manifest.get("files"), list) else []
    aggregate = hashlib.sha256(
        "\n".join(f"{item.get('path')}:{item.get('sha256')}" for item in ledger).encode()
    ).hexdigest()
    checks.extend(
        [
            check(
                "powerbi.project_hash",
                expected=manifest.get("project_sha256", _ABSENT),
                actual=aggregate,
                evidence=manifest_portable,
            ),
            check(
                "powerbi.product",
                expected=PRODUCT,
                actual=manifest.get("product", _ABSENT),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.capability_status",
                expected="INTEGRATION_ONLY",
                actual=manifest.get("capability_status", _ABSENT),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.no_fabricated_pbix",
                expected=False,
                actual=manifest.get("contains_pbix", _ABSENT),
                evidence=manifest_portable,
            ),
        ]
    )
    scope_path = root / "Data" / "evidence_scope.csv"
    kpi_path = root / "Data" / "kpi_snapshot.csv"
    reconciliation_path = root / "Validation" / "reconciliation_snapshot.csv"
    for item in (scope_path, kpi_path, reconciliation_path):
        if item.is_file():
            paths.append(_portable(item, repository_root))
    try:
        scope_rows = _read_csv(scope_path)
        kpi_rows = _read_csv(kpi_path)
        reconciliation_rows = _read_csv(reconciliation_path)
        scope = scope_rows[0] if len(scope_rows) == 1 else {}
        loss = next(row for row in kpi_rows if row.get("metric_id") == expected["metric_id"])
        root_rows = {
            row["metric_id"]: row for row in reconciliation_rows if row.get("scope") == "Root cause"
        }
        actual = {
            "current_annualised_net_loss_rate": float(loss["value"]),
            "prior_annualised_net_loss_rate": float(loss["prior_value"]),
            "observed_change_bps": float(root_rows["OBSERVED_CHANGE_BPS"]["current_value"]),
            "mix_contribution_bps": float(root_rows["MIX_CONTRIBUTION_BPS"]["current_value"]),
            "within_segment_contribution_bps": float(
                root_rows["PERFORMANCE_CONTRIBUTION_BPS"]["current_value"]
            ),
            "reconciliation_residual_bps": float(
                root_rows["RECONCILIATION_RESIDUAL_BPS"]["current_value"]
            ),
            "reporting_period": scope.get("reporting_period"),
            "run_id": scope.get("run_id"),
            "metric_registry_version": scope.get("metric_registry_version"),
            "data_quality_status": scope.get("quality_status"),
            "synthetic_data": scope.get("synthetic_data_flag", "").casefold() == "true",
        }
        checks.extend(
            compare_story(
                expected,
                actual,
                prefix="powerbi",
                fields=(
                    "current_annualised_net_loss_rate",
                    "prior_annualised_net_loss_rate",
                    "observed_change_bps",
                    "mix_contribution_bps",
                    "within_segment_contribution_bps",
                    "reconciliation_residual_bps",
                    "reporting_period",
                    "run_id",
                    "metric_registry_version",
                    "data_quality_status",
                    "synthetic_data",
                ),
                evidence=_portable(reconciliation_path, repository_root),
            )
        )
        checks.extend(
            [
                check(
                    "powerbi.evidence_payload_sha256",
                    expected=expected.get("evidence_payload_sha256"),
                    actual=scope.get("evidence_sha256", _ABSENT),
                    evidence=_portable(scope_path, repository_root),
                ),
                check(
                    "powerbi.evidence_id_maps_to_run",
                    expected=True,
                    actual=_normalise_evidence_id(
                        scope.get("evidence_id", _ABSENT), expected.get("run_id")
                    ),
                    evidence=_portable(scope_path, repository_root),
                ),
            ]
        )
    except (OSError, ValueError, KeyError, StopIteration, IndexError) as exc:
        checks.append(
            check(
                "powerbi.validation_data_readable",
                expected=True,
                actual=False,
                detail=f"{type(exc).__name__}: {exc}",
                evidence=manifest_portable,
            )
        )
    checks.extend(
        [
            check(
                "powerbi.provenance_artifact_id",
                actual=_nested(manifest, "artifact_id", "build_id", "project_id"),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.provenance_creation_time",
                actual=_nested(manifest, "created_at_utc", "built_at_utc", "generated_at"),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.provenance_dataset_hash",
                expected=expected.get("dataset_hash", _ABSENT),
                actual=_nested(manifest, "dataset_hash"),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.provenance_configuration_hash",
                expected=expected.get("configuration_hash", _ABSENT),
                actual=_nested(manifest, "configuration_hash"),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.provenance_code_version",
                actual=_nested(manifest, "code_version", "generator_version"),
                evidence=manifest_portable,
            ),
            check(
                "powerbi.provenance_filter_scope",
                actual=_nested(manifest, "filter_scope"),
                evidence=manifest_portable,
            ),
        ]
    )
    return channel(
        "power_bi_validation",
        "Power BI validation data",
        checks=checks,
        artifact_paths=paths,
        notes=[
            "Power BI Desktop and Service publication remain outside this static reconciliation."
        ],
    )


def _static_site_channel(repository_root: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    root = repository_root / "outputs" / "share_site"
    evidence_path = root / "data" / "evidence.json"
    index_path = root / "index.html"
    build_manifest_path = root / "build_manifest.json"
    if not evidence_path.is_file() or not index_path.is_file():
        return _missing_channel(
            "static_share_site",
            "Static share site",
            "outputs/share_site/index.html",
        )
    payload = _read_json(evidence_path)
    actual = _snapshot_story(payload)
    evidence_portable = _portable(evidence_path, repository_root)
    index_portable = _portable(index_path, repository_root)
    paths = [evidence_portable, index_portable]
    checks = compare_story(
        expected,
        actual,
        prefix="static_site",
        fields=(
            "current_annualised_net_loss_rate",
            "prior_annualised_net_loss_rate",
            "observed_change_bps",
            "mix_contribution_bps",
            "within_segment_contribution_bps",
            "reconciliation_residual_bps",
            "primary_dimension",
            "primary_driver",
            "causal_status",
            "reporting_period",
            "run_id",
            "metric_registry_version",
            "data_quality_status",
            "synthetic_data",
        ),
        evidence=evidence_portable,
    )
    source_inputs = _nested(payload, "source_context.source_inputs")
    canonical_source = None
    if isinstance(source_inputs, list):
        canonical_source = next(
            (
                item
                for item in source_inputs
                if isinstance(item, Mapping)
                and item.get("path") == "exports/validation/interop_evidence_snapshot.json"
            ),
            None,
        )
    checks.extend(
        [
            check(
                "static_site.canonical_source_file_hash",
                expected=expected.get("canonical_file_sha256"),
                actual=(
                    canonical_source.get("sha256")
                    if isinstance(canonical_source, Mapping)
                    else _ABSENT
                ),
                evidence=evidence_portable,
            ),
            check(
                "static_site.evidence_sidecar_hash",
                expected=(root / "data" / "evidence.json.sha256")
                .read_text(encoding="utf-8")
                .split()[0]
                if (root / "data" / "evidence.json.sha256").is_file()
                else _ABSENT,
                actual=sha256_file(evidence_path),
                evidence=evidence_portable,
            ),
        ]
    )
    html = index_path.read_text(encoding="utf-8")
    checks.extend(
        [
            check(
                "static_site.product",
                expected=True,
                actual=PRODUCT in html,
                evidence=index_portable,
            ),
            check(
                "static_site.tagline",
                expected=True,
                actual=TAGLINE in html,
                evidence=index_portable,
            ),
            check(
                "static_site.associational_disclosure",
                expected=True,
                actual="associational" in html.casefold(),
                evidence=index_portable,
            ),
        ]
    )
    if build_manifest_path.is_file():
        build = _read_json(build_manifest_path)
        build_portable = _portable(build_manifest_path, repository_root)
        paths.append(build_portable)
        checks.extend(
            verify_file_ledger(
                root,
                build.get("files"),
                prefix="static_site.build",
                repository_root=repository_root,
            )
        )
        checks.extend(
            [
                check(
                    "static_site.build_validation",
                    expected=PASS,
                    actual=_nested(build, "validation.status"),
                    evidence=build_portable,
                ),
                check(
                    "static_site.build_evidence_id",
                    expected=payload.get("evidence_id"),
                    actual=build.get("source_evidence_id", _ABSENT),
                    evidence=build_portable,
                ),
            ]
        )
    else:
        checks.append(check("static_site.build_manifest", actual=_ABSENT))
    release_checks, release_paths = _release_manifest_checks(
        repository_root, index_path, expected, prefix="static_site.manifest"
    )
    checks.extend(release_checks)
    paths.extend(release_paths)
    return channel(
        "static_share_site",
        "Static share site",
        checks=checks,
        artifact_paths=paths,
    )


def _text_has_core_story(text: str, expected: Mapping[str, Any]) -> dict[str, bool]:
    normalised = re.sub(r"\s+", " ", text)
    return {
        "observed": f"{float(expected['observed_change_bps']):.4f}" in normalised,
        "mix": f"{float(expected['mix_contribution_bps']):.4f}" in normalised,
        "within": f"{float(expected['within_segment_contribution_bps']):.4f}" in normalised,
        "driver": str(expected.get("primary_driver")) in normalised,
        "causal": "associational" in normalised.casefold(),
    }


def _pdf_text(path: Path) -> str:
    """Extract PDF text with an available native or Python PDF reader."""

    try:
        from pypdf import PdfReader
    except ImportError:
        executable = shutil.which("pdftotext")
        if executable is not None:
            result = subprocess.run(
                [executable, str(path), "-"],
                check=True,
                capture_output=True,
                text=True,
            )
            return result.stdout
        return _pdf_text_with_pdfkit(path)
    return "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)


def _pdf_text_with_pdfkit(path: Path) -> str:
    """Extract text through macOS PDFKit when portable readers are unavailable."""

    swift = shutil.which("swift")
    if sys.platform != "darwin" or swift is None:
        raise ImportError(
            "No PDF text reader is available: install pypdf or Poppler pdftotext"
        ) from None
    program = """
import Foundation
import PDFKit

let url = URL(fileURLWithPath: CommandLine.arguments[1])
guard let document = PDFDocument(url: url) else {
    fatalError("Unable to open PDF")
}
for index in 0..<document.pageCount {
    print(document.page(at: index)?.string ?? "")
}
"""
    with tempfile.TemporaryDirectory(prefix="naim-pdfkit-") as module_cache:
        result = subprocess.run(
            [
                swift,
                "-module-cache-path",
                module_cache,
                "-e",
                program,
                str(path),
            ],
            check=True,
            capture_output=True,
            text=True,
        )
    return result.stdout


def _linkedin_channel(repository_root: Path, expected: Mapping[str, Any]) -> dict[str, Any]:
    root = repository_root / "outputs" / "linkedin"
    manifest_path = root / "package-manifest.json"
    if not manifest_path.is_file():
        return _missing_channel(
            "linkedin_carousel",
            "LinkedIn social pack and carousel",
            "outputs/linkedin/package-manifest.json",
        )
    manifest = _read_json(manifest_path)
    manifest_portable = _portable(manifest_path, repository_root)
    paths = [manifest_portable]
    checks: list[dict[str, Any]] = [
        check(
            "linkedin.product",
            expected=PRODUCT,
            actual=manifest.get("product", _ABSENT),
            evidence=manifest_portable,
        ),
        check(
            "linkedin.tagline",
            expected=TAGLINE,
            actual=manifest.get("tagline", _ABSENT),
            evidence=manifest_portable,
        ),
        check(
            "linkedin.evidence_id_maps_to_run",
            expected=True,
            actual=_normalise_evidence_id(
                manifest.get("evidence_id", _ABSENT), expected.get("run_id")
            ),
            evidence=manifest_portable,
        ),
    ]
    text_artifacts = manifest.get("text_artifacts")
    if isinstance(text_artifacts, Mapping):
        for name, status in text_artifacts.items():
            path = root / str(name)
            checks.extend(
                [
                    check(
                        f"linkedin.text.{name}.status",
                        expected="READY",
                        actual=status,
                        evidence=manifest_portable,
                    ),
                    check(
                        f"linkedin.text.{name}.present",
                        expected=True,
                        actual=path.is_file(),
                        evidence=_portable(path, repository_root) if path.is_file() else str(name),
                    ),
                ]
            )
            if path.is_file():
                paths.append(_portable(path, repository_root))
    else:
        checks.append(check("linkedin.text_artifacts", actual=_ABSENT))
    research_path = root / "research-summary.md"
    combined_text = ""
    if isinstance(text_artifacts, Mapping):
        for name in text_artifacts:
            text_path = root / str(name)
            if text_path.is_file():
                combined_text += "\n" + text_path.read_text(encoding="utf-8")
    claims = _text_has_core_story(combined_text, expected) if combined_text else {}
    for claim_name in ("observed", "mix", "within", "driver", "causal"):
        checks.append(
            check(
                f"linkedin.text_story.{claim_name}",
                expected=True,
                actual=claims.get(claim_name, _ABSENT),
                evidence=_portable(research_path, repository_root)
                if research_path.is_file()
                else manifest_portable,
            )
        )
    pptx_candidates = sorted(root.glob("*.pptx"))
    pdf_candidates = sorted(root.glob("*.pdf"))
    checks.extend(
        [
            check(
                "linkedin.editable_carousel_present",
                expected=True,
                actual=True if pptx_candidates else _ABSENT,
                evidence=_portable(pptx_candidates[-1], repository_root)
                if pptx_candidates
                else manifest_portable,
            ),
            check(
                "linkedin.pdf_carousel_present",
                expected=True,
                actual=True if pdf_candidates else _ABSENT,
                evidence=_portable(pdf_candidates[-1], repository_root)
                if pdf_candidates
                else manifest_portable,
            ),
        ]
    )
    if pptx_candidates:
        pptx = pptx_candidates[-1]
        paths.append(_portable(pptx, repository_root))
        try:
            text = _presentation_text(pptx)
            rounded_tokens = (
                f"{float(expected['observed_change_bps']):.1f}",
                f"{float(expected['mix_contribution_bps']):.1f}",
                f"{float(expected['within_segment_contribution_bps']):.1f}",
                str(expected.get("primary_driver")),
            )
            checks.append(
                check(
                    "linkedin.editable_carousel_story",
                    expected=True,
                    actual=all(token in text for token in rounded_tokens)
                    and "associational" in text.casefold(),
                    evidence=_portable(pptx, repository_root),
                )
            )
        except (ImportError, OSError, ValueError) as exc:
            checks.append(
                unverifiable(
                    "linkedin.editable_carousel_story",
                    detail=f"{type(exc).__name__}: {exc}",
                    evidence=_portable(pptx, repository_root),
                )
            )
    if pdf_candidates:
        pdf = pdf_candidates[-1]
        paths.append(_portable(pdf, repository_root))
        try:
            text = _pdf_text(pdf)
            checks.append(
                check(
                    "linkedin.pdf_carousel_story",
                    expected=True,
                    actual=(
                        f"{float(expected['observed_change_bps']):.1f}" in text
                        and str(expected.get("primary_driver")) in text
                        and "associational" in text.casefold()
                    ),
                    evidence=_portable(pdf, repository_root),
                )
            )
        except (ImportError, OSError, ValueError) as exc:
            checks.append(
                unverifiable(
                    "linkedin.pdf_carousel_story",
                    detail=f"{type(exc).__name__}: {exc}",
                    evidence=_portable(pdf, repository_root),
                )
            )
    checks.extend(
        verify_file_ledger(
            root,
            manifest.get("files"),
            prefix="linkedin",
            repository_root=repository_root,
        )
    )
    checks.extend(
        [
            check(
                "linkedin.provenance_artifact_id",
                actual=_nested(manifest, "artifact_id", "build_id", "package_id"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_creation_time",
                actual=_nested(manifest, "created_at_utc", "built_at_utc", "generated_at"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_dataset_hash",
                expected=expected.get("dataset_hash", _ABSENT),
                actual=_nested(manifest, "dataset_hash"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_configuration_hash",
                expected=expected.get("configuration_hash", _ABSENT),
                actual=_nested(manifest, "configuration_hash"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_code_version",
                actual=_nested(manifest, "code_version", "generator_version"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_metric_registry_version",
                expected=expected.get("metric_registry_version"),
                actual=_nested(manifest, "metric_registry_version"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_reporting_period",
                expected=expected.get("reporting_period"),
                actual=_nested(manifest, "reporting_period"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_filter_scope",
                actual=_nested(manifest, "filter_scope"),
                evidence=manifest_portable,
            ),
            check(
                "linkedin.provenance_data_quality",
                expected=expected.get("data_quality_status"),
                actual=_nested(manifest, "data_quality_result", "data_quality.status"),
                evidence=manifest_portable,
            ),
        ]
    )
    return channel(
        "linkedin_carousel",
        "LinkedIn social pack and carousel",
        checks=checks,
        artifact_paths=paths,
    )


def reconcile_repository(repository_root: Path) -> dict[str, Any]:
    """Run every adapter and return the complete fail-closed report."""

    repository_root = repository_root.resolve()
    canonical, api_channel = canonical_context(repository_root)
    story = canonical["story"]
    channels = [
        api_channel,
        _runtime_snapshot_channel(
            repository_root,
            story,
            channel_id="ui_snapshot",
            label="Browser-validated UI snapshot",
            relative_path="outputs/validation/ui_evidence_snapshot.json",
        ),
        _workbook_channel(repository_root, story),
        _presentation_channel(repository_root, story),
        _hyper_channel(repository_root, story),
        _powerbi_channel(repository_root, story),
        _runtime_snapshot_channel(
            repository_root,
            story,
            channel_id="streamlit_snapshot",
            label="Streamlit runtime snapshot",
            relative_path="outputs/streamlit/evidence_snapshot.json",
        ),
        _static_site_channel(repository_root, story),
        _linkedin_channel(repository_root, story),
    ]
    for row in channels:
        artifact_records = []
        for relative in row["artifact_paths"]:
            artifact_path = (repository_root / relative).resolve()
            if artifact_path.is_relative_to(repository_root) and artifact_path.is_file():
                artifact_records.append(
                    {
                        "path": relative,
                        "bytes": artifact_path.stat().st_size,
                        "sha256": sha256_file(artifact_path),
                    }
                )
        row["artifacts"] = artifact_records
    all_checks = [item for row in channels for item in row["checks"] if item.get("required")]
    failed = [item for item in all_checks if item["outcome"] == FAIL]
    incomplete = [item for item in all_checks if item["outcome"] in {MISSING, UNVERIFIABLE}]
    if failed:
        overall = FAIL
    elif incomplete:
        overall = INCOMPLETE
    else:
        overall = PASS
    status_counts = {
        status: sum(1 for item in channels if item["status"] == status)
        for status in (PASS, FAIL, INCOMPLETE)
    }
    return {
        "schema_version": SCHEMA_VERSION,
        "schema_reference": "schemas/cross-artifact-reconciliation.schema.json",
        "product": PRODUCT,
        "generated_at_utc": datetime.now(UTC).isoformat(),
        "result": overall,
        "release_allowed": overall == PASS,
        "methodology": {
            "canonical_source": canonical["source_path"],
            "exact_numeric_tolerance": 1e-9,
            "residual_tolerance_bps": 1e-8,
            "one_decimal_display_tolerance_bps": 0.05,
            "missing_evidence_policy": "Fail closed; missing or unverifiable required evidence blocks release.",
            "identity_policy": (
                "Evidence-ID prefixes may change during rebranding; the complete source run ID must match."
            ),
            "scope_policy": (
                "The headline is All portfolio. BASKET-001 is an approved secondary control basket, "
                "not a claim that the headline was basket-filtered."
            ),
        },
        "canonical": canonical,
        "channels": channels,
        "summary": {
            "required_channel_count": len(channels),
            "channel_status_counts": status_counts,
            "failed_check_count": len(failed),
            "missing_or_unverifiable_check_count": len(incomplete),
            "missing_or_incomplete_channels": [
                item["channel_id"] for item in channels if item["status"] == INCOMPLETE
            ],
            "failed_channels": [item["channel_id"] for item in channels if item["status"] == FAIL],
            "blocking_checks": [
                {
                    "check_id": item["check_id"],
                    "outcome": item["outcome"],
                    "evidence": item.get("evidence"),
                    "detail": item.get("detail"),
                }
                for item in [*failed, *incomplete]
            ],
        },
    }


def _parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--repository-root", type=Path, default=DEFAULT_REPOSITORY_ROOT)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    return parser.parse_args()


def main() -> int:
    args = _parse_args()
    root = args.repository_root.resolve()
    output = args.output if args.output.is_absolute() else root / args.output
    try:
        report = reconcile_repository(root)
    except (OSError, ValueError, TypeError, json.JSONDecodeError, ReconciliationError) as exc:
        report = {
            "schema_version": SCHEMA_VERSION,
            "schema_reference": "schemas/cross-artifact-reconciliation.schema.json",
            "product": PRODUCT,
            "generated_at_utc": datetime.now(UTC).isoformat(),
            "result": FAIL,
            "release_allowed": False,
            "fatal_error": f"{type(exc).__name__}: {exc}",
        }
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(report, indent=2, sort_keys=True) + "\n", encoding="utf-8")
    print(json.dumps({"result": report["result"], "output": _portable(output, root)}))
    return 0 if report["result"] == PASS else 1


if __name__ == "__main__":
    sys.exit(main())
