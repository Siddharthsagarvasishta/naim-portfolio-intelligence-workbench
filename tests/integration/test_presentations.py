from __future__ import annotations

import io
import json
import re
import zipfile
from pathlib import Path

import pytest
from fastapi.testclient import TestClient
from pptx import Presentation

import naim_risk.api as api_main
from naim_risk.api import app, get_presentation_output_root, get_service
from naim_risk.config import load_config
from naim_risk.presentations import (
    generate_presentation,
    presentation_record,
    resolve_presentation_file,
    validate_presentation,
)
from naim_risk.runtime_modes import DataMode, source_context
from naim_risk.service import WorkbenchService
from naim_risk.workflow import WorkflowStore


def test_editable_presentation_has_notes_manifest_and_safe_persistence(tmp_path: Path) -> None:
    config = load_config("test", data_root=tmp_path / "data")
    service = WorkbenchService(config)
    store = WorkflowStore(f"sqlite+pysqlite:///{(tmp_path / 'workflow.sqlite3').resolve()}")
    context = source_context(config, DataMode.DEMO)
    output_root = tmp_path / "outputs" / "presentations"
    payload = {
        "reporting_period": service.metadata()["as_of"],
        "comparison_period": "2024-07-01",
        "workspace_id": service.workspaces()["data"][0]["workspace_id"],
        "scenario_name": "Baseline",
        "selected_sections": [
            "executive_summary",
            "kpis",
            "root_cause",
            "partners",
            "decision_log",
        ],
        "include_appendix": True,
    }
    result = generate_presentation(
        service,
        payload,
        store=store,
        source_context=context,
        actor="portfolio.analyst",
        output_root=output_root,
    )
    assert "path" not in result
    assert result["filename"] == "nAIM_Portfolio_Intelligence_Review_2024_08.pptx"
    assert result["slide_count"] == 7
    assert result["validation_status"] == "PASS"
    assert result["approval_state"] == "DRAFT"

    path = resolve_presentation_file(store, result["presentation_id"], output_root=output_root)
    manifest_path = resolve_presentation_file(
        store,
        result["presentation_id"],
        manifest=True,
        output_root=output_root,
    )
    validation = validate_presentation(
        path,
        [
            "executive_summary",
            "kpis",
            "root_cause",
            "partners",
            "decision_log",
            "appendix",
        ],
    )
    assert validation["status"] == "PASS"
    deck = Presentation(path)
    assert all(slide.notes_slide.notes_text_frame.text.strip() for slide in deck.slides)
    assert any(shape.has_chart for slide in deck.slides for shape in slide.shapes)
    assert any(shape.has_table for slide in deck.slides for shape in slide.shapes)
    with zipfile.ZipFile(path) as archive:
        chart_xml = b"".join(
            archive.read(name)
            for name in archive.namelist()
            if name.startswith("ppt/charts/") and name.endswith(".xml")
        )
    assert not re.search(rb'<c:(?:axId|crossAx)\s+val="-', chart_xml)

    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    assert manifest["artifact"]["filename"] == result["filename"]
    assert manifest["artifact"]["sha256"]
    assert manifest["validation"]["status"] == "PASS"
    assert manifest["evidence_ids"]
    assert manifest["reconciliation"]["status"] == "API_RECONCILED"
    assert manifest["filter_scope"]["headline_scope"] == "all_portfolio"
    assert manifest["filter_scope"]["approved_reference_basket"] is None

    record = presentation_record(store, result["presentation_id"])
    assert record["approval_state"] == "DRAFT"
    store.close()


@pytest.mark.integration
def test_live_presentation_api_generates_and_downloads_without_server_path(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    config = load_config("test", data_root=tmp_path / "data")
    service = WorkbenchService(config)
    output_root = tmp_path / "outputs" / "presentations"
    monkeypatch.setenv("NAIM_AUTH_MODE", "disabled")
    monkeypatch.setenv("NAIM_DATA_MODE", "DEMO")
    monkeypatch.setenv("NAIM_DATASET_PROFILE", "test")
    monkeypatch.setenv("NAIM_DATA_DIR", str(tmp_path / "data"))
    monkeypatch.setenv(
        "NAIM_DATABASE_URL",
        f"sqlite+pysqlite:///{(tmp_path / 'api-workflow.sqlite3').resolve()}",
    )
    api_main.reset_application_state()
    app.dependency_overrides[get_service] = lambda: service
    app.dependency_overrides[get_presentation_output_root] = lambda: output_root
    try:
        client = TestClient(app)
        response = client.post(
            "/api/v1/presentations/generate",
            json={
                "reporting_period": service.metadata()["as_of"],
                "selected_sections": ["executive_summary", "kpis", "root_cause"],
                "include_appendix": True,
                "speaker_notes": True,
            },
        )
        assert response.status_code == 201
        payload = response.json()
        assert "path" not in payload
        assert payload["data_mode"] == "DEMO"
        presentation_id = payload["presentation_id"]

        status = client.get(f"/api/v1/presentations/{presentation_id}")
        assert status.status_code == 200
        assert status.json()["validation_status"] == "PASS"
        download = client.get(payload["download_url"])
        assert download.status_code == 200
        deck = Presentation(io.BytesIO(download.content))
        assert len(deck.slides) == 5
        manifest = client.get(payload["manifest_url"])
        assert manifest.status_code == 200
        assert manifest.json()["presentation_id"] == presentation_id
    finally:
        app.dependency_overrides.clear()
        api_main.reset_application_state()
