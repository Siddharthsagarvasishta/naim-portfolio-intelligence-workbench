"""Portable, reproducible exports for Market Risk and Volatility Lab evidence."""

from __future__ import annotations

import hashlib
import json
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

import numpy as np
import pandas as pd

from naim_risk.market_risk.providers import MarketDataError, MarketPriceFrame


def _json_default(value: Any) -> Any:
    if isinstance(value, (np.integer,)):
        return int(value)
    if isinstance(value, (np.floating,)):
        return None if not np.isfinite(value) else float(value)
    if isinstance(value, (np.bool_,)):
        return bool(value)
    if isinstance(value, (pd.Timestamp, datetime)):
        return value.isoformat()
    raise TypeError(f"Unsupported JSON value: {type(value).__name__}")


def _json_text(value: Any, *, indent: int | None = None) -> str:
    return json.dumps(
        value,
        default=_json_default,
        allow_nan=False,
        indent=indent,
        sort_keys=True,
    )


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1_048_576), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _records(value: Any) -> pd.DataFrame:
    if isinstance(value, list):
        return pd.json_normalize(value) if value else pd.DataFrame([{"status": "no rows"}])
    if isinstance(value, dict):
        return pd.json_normalize(value, sep=".")
    return pd.DataFrame([{"value": value}])


def _flatten_key_values(value: dict[str, Any], prefix: str = "") -> pd.DataFrame:
    rows: list[dict[str, Any]] = []

    def visit(item: Any, path: str) -> None:
        if isinstance(item, dict):
            for key, nested in item.items():
                visit(nested, f"{path}.{key}" if path else str(key))
        elif isinstance(item, list):
            rows.append({"field": path, "value": _json_text(item)})
        else:
            rows.append({"field": path, "value": item})

    visit(value, prefix)
    return pd.DataFrame(rows)


def _write_excel(
    path: Path,
    analysis: dict[str, Any],
    market: MarketPriceFrame | None,
) -> None:
    historical_rows = []
    historical = analysis.get("historical_volatility", {})
    for name, details in historical.get("estimators", {}).items():
        historical_rows.append({"estimator": name, **details})
    comparison = pd.DataFrame(analysis.get("model_comparison", {}).get("models", []))
    parameters = []
    for model, details in analysis.get("conditional_volatility", {}).items():
        for parameter in details.get("parameters", []):
            parameters.append({"model": model, **parameter})
    overview = {
        "module": analysis.get("module"),
        "status": analysis.get("status"),
        "purpose": analysis.get("purpose"),
        "instrument": analysis.get("source", {}).get("instrument"),
        "provider": analysis.get("source", {}).get("provider"),
        "price_basis": analysis.get("source", {}).get("price_basis"),
    }
    assumptions = {
        "source_notes": analysis.get("source", {}).get("notes", []),
        "implied_volatility": analysis.get("implied_volatility", {}).get("assumptions", []),
        "governance": analysis.get("governance", {}).get("limitations", []),
    }
    methodology = {
        "return_method": analysis.get("returns", {}).get("selected_return"),
        "frequency": analysis.get("returns", {}).get("frequency"),
        "annualisation_factor": analysis.get("returns", {}).get("annualisation_factor"),
        "historical_estimators": list(historical.get("estimators", {})),
        "conditional_models": list(analysis.get("conditional_volatility", {})),
        "var_warning": analysis.get("var_expected_shortfall", {}).get("warning"),
    }
    refresh = {
        "generated_at_utc": datetime.now(UTC).isoformat(),
        "schema_version": analysis.get("schema_version"),
        "source_retrieval_time": analysis.get("source", {}).get("retrieval_time"),
        "source_sha256": analysis.get("source", {}).get("raw_source_sha256"),
        "source_is_synthetic": analysis.get("source", {}).get("source_is_synthetic"),
        "external_data_used": analysis.get("governance", {}).get("external_data_used"),
    }
    sheets = {
        "Overview": _flatten_key_values(overview),
        "Price Data": market.data.copy()
        if market is not None
        else pd.DataFrame([{"status": "not supplied"}]),
        "Return Statistics": _flatten_key_values(analysis.get("returns", {}).get("summary", {})),
        "Historical Volatility": pd.DataFrame(historical_rows),
        "EWMA": _flatten_key_values(analysis.get("ewma", {})),
        "ARCH GARCH": pd.DataFrame(parameters),
        "Diagnostics": _flatten_key_values(analysis.get("diagnostics", {})),
        "Implied Volatility": _flatten_key_values(analysis.get("implied_volatility", {})),
        "VaR Expected Shortfall": _flatten_key_values(analysis.get("var_expected_shortfall", {})),
        "Backtesting": _flatten_key_values(analysis.get("var_backtesting", {})),
        "Model Comparison": comparison,
        "Assumptions": _flatten_key_values(assumptions),
        "Methodology": _flatten_key_values(methodology),
        "Refresh Control": _flatten_key_values(refresh),
    }
    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        for name, frame in sheets.items():
            safe = frame.copy()
            for column in safe.columns:
                safe[column] = safe[column].map(
                    lambda value: (
                        _json_text(value)
                        if isinstance(value, (dict, list))
                        else value.isoformat()
                        if isinstance(value, (pd.Timestamp, datetime))
                        else value
                    )
                )
            safe.to_excel(writer, sheet_name=name, index=False)
        workbook = writer.book
        workbook.properties.title = "nAIM Market Risk and Volatility Lab"
        workbook.properties.subject = "Reproducible market-risk model evidence"
        workbook.properties.creator = "nAIM Portfolio Intelligence Workbench"
        for worksheet in workbook.worksheets:
            worksheet.freeze_panes = "A2"
            worksheet.auto_filter.ref = worksheet.dimensions
            for cell in worksheet[1]:
                cell.font = cell.font.copy(bold=True, color="FFFFFF")
                cell.fill = cell.fill.copy(fill_type="solid", fgColor="17324D")
            for column in worksheet.columns:
                letter = column[0].column_letter
                maximum = max(len(str(cell.value or "")) for cell in column)
                worksheet.column_dimensions[letter].width = min(max(maximum + 2, 12), 45)


def _add_text_slide(presentation: Any, title: str, lines: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title
    frame = slide.placeholders[1].text_frame
    frame.clear()
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(line)
        paragraph.level = 0


def _write_presentation(path: Path, analysis: dict[str, Any]) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except (ImportError, ModuleNotFoundError) as exc:
        raise MarketDataError(
            f"python-pptx is required for editable presentation export: {exc}"
        ) from exc
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = "nAIM Market Risk and Volatility Lab"
    title_slide.placeholders[1].text = (
        "Name the movement. Own the evidence.\n"
        "Quantitative diagnostics — not a trading recommendation"
    )
    source = analysis.get("source", {})
    summary = analysis.get("returns", {}).get("summary", {})
    _add_text_slide(
        presentation,
        "Source and return preparation",
        [
            f"Instrument: {source.get('instrument', 'N/A')}",
            f"Provider: {source.get('provider', 'N/A')}",
            f"Price basis: {source.get('price_basis', 'N/A')}",
            f"Observations: {summary.get('observations', 'N/A')}",
            f"Annualised volatility: {summary.get('annualised_standard_deviation', 'N/A')}",
            f"Outliers: {summary.get('outliers', 'N/A')}",
        ],
    )
    model_lines = []
    for row in analysis.get("model_comparison", {}).get("models", []):
        model_lines.append(
            f"{row['model']}: latest vol {row.get('latest_annualised_volatility')}; "
            f"held-out QLIKE {row.get('out_of_sample_qlike')}"
        )
    _add_text_slide(presentation, "Model comparison", model_lines or ["No model rows available"])
    var_methods = analysis.get("var_expected_shortfall", {}).get("methods", {})
    _add_text_slide(
        presentation,
        "Tail risk and governance",
        [
            *[
                f"{name}: VaR {details.get('var')}; ES {details.get('expected_shortfall')}"
                for name, details in var_methods.items()
            ],
            "VaR is not the maximum possible loss.",
            "External risk-regime overlay is associational only.",
            "No output is investment advice or a trade signal.",
        ],
    )
    presentation.core_properties.title = "nAIM Market Risk and Volatility Lab"
    presentation.core_properties.author = "nAIM Portfolio Intelligence Workbench"
    presentation.core_properties.subject = "Editable market-risk evidence"
    presentation.save(path)


def _write_notebook(path: Path) -> None:
    notebook = {
        "cells": [
            {
                "cell_type": "markdown",
                "metadata": {},
                "source": [
                    "# nAIM Market Risk and Volatility Lab reproduction\n",
                    "This notebook reloads exported evidence. It is not a trading recommendation.\n",
                ],
            },
            {
                "cell_type": "code",
                "execution_count": None,
                "metadata": {},
                "outputs": [],
                "source": [
                    "import json\n",
                    "from pathlib import Path\n",
                    "evidence = json.loads(Path('market_risk_evidence.json').read_text())\n",
                    "evidence['source'], evidence['model_comparison']['models']\n",
                ],
            },
        ],
        "metadata": {
            "kernelspec": {"display_name": "Python 3", "language": "python", "name": "python3"},
            "language_info": {"name": "python", "version": "3.12"},
        },
        "nbformat": 4,
        "nbformat_minor": 5,
    }
    path.write_text(_json_text(notebook, indent=2) + "\n", encoding="utf-8")


def export_market_risk_bundle(
    analysis: dict[str, Any],
    destination: str | Path,
    *,
    market: MarketPriceFrame | None = None,
    include_excel: bool = True,
    include_presentation: bool = True,
) -> dict[str, Any]:
    """Write portable evidence files and return a path-safe integrity manifest."""

    if analysis.get("module") != "Market Risk and Volatility Lab":
        raise MarketDataError("Export requires Market Risk and Volatility Lab evidence")
    root = Path(destination).expanduser().resolve()
    if root == Path(root.anchor):
        raise MarketDataError("A filesystem root cannot be used as an export destination")
    root.mkdir(parents=True, exist_ok=True)
    written: list[Path] = []
    evidence_path = root / "market_risk_evidence.json"
    evidence_path.write_text(_json_text(analysis, indent=2) + "\n", encoding="utf-8")
    written.append(evidence_path)
    returns = pd.DataFrame(analysis.get("returns", {}).get("observations", []))
    returns_csv = root / "prepared_returns.csv"
    returns.to_csv(returns_csv, index=False)
    written.append(returns_csv)
    returns_parquet = root / "prepared_returns.parquet"
    returns.to_parquet(returns_parquet, index=False)
    written.append(returns_parquet)
    if market is not None:
        raw_csv = root / "price_data.csv"
        market.data.to_csv(raw_csv, index=False)
        written.append(raw_csv)
        raw_parquet = root / "price_data.parquet"
        market.data.to_parquet(raw_parquet, index=False)
        written.append(raw_parquet)
    chart_data = {
        "rolling_volatility": analysis.get("historical_volatility", {}).get("rolling", {}),
        "ewma": analysis.get("ewma", {}).get("series", []),
        "regimes": analysis.get("regimes", {}).get("series", []),
        "model_comparison": analysis.get("model_comparison", {}).get("models", []),
    }
    chart_path = root / "chart_data.json"
    chart_path.write_text(_json_text(chart_data, indent=2) + "\n", encoding="utf-8")
    written.append(chart_path)
    notebook_path = root / "reproduce_market_risk.ipynb"
    _write_notebook(notebook_path)
    written.append(notebook_path)
    if include_excel:
        workbook_path = root / "nAIM_Market_Risk_Volatility_Lab.xlsx"
        _write_excel(workbook_path, analysis, market)
        written.append(workbook_path)
    if include_presentation:
        presentation_path = root / "nAIM_Market_Risk_Volatility_Lab.pptx"
        _write_presentation(presentation_path, analysis)
        written.append(presentation_path)
    manifest = {
        "schema_version": "1.0.0",
        "module": "Market Risk and Volatility Lab",
        "created_at_utc": datetime.now(UTC).isoformat(),
        "source_sha256": analysis.get("source", {}).get("raw_source_sha256"),
        "files": [
            {
                "path": path.relative_to(root).as_posix(),
                "bytes": path.stat().st_size,
                "sha256": _sha256(path),
            }
            for path in sorted(written)
        ],
        "capabilities": {
            "csv": "generated",
            "parquet": "generated",
            "json_evidence": "generated",
            "notebook": "generated",
            "chart_data": "generated",
            "excel": "generated" if include_excel else "not_requested",
            "editable_powerpoint": "generated" if include_presentation else "not_requested",
            "pdf": "not_generated_requires_separate_render_and_visual_validation",
        },
        "warnings": [
            "No server filesystem paths are included in this manifest.",
            "Source redistribution rights remain the user's responsibility.",
            "Outputs are analytical evidence, not trading recommendations.",
        ],
    }
    manifest_path = root / "export_manifest.json"
    manifest_path.write_text(_json_text(manifest, indent=2) + "\n", encoding="utf-8")
    return manifest
