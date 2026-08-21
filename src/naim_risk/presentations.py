"""Editable, evidence-linked PowerPoint generation for the nAIM workbench."""

from __future__ import annotations

import hashlib
import json
import os
import re
import tempfile
import zipfile
from collections.abc import Mapping, Sequence
from datetime import UTC, datetime
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from naim_risk.config import REPOSITORY_ROOT
from naim_risk.runtime_modes import SourceContext
from naim_risk.service import WorkbenchService
from naim_risk.workflow import DuplicateObject, WorkflowStore

PRODUCT = "nAIM Portfolio Intelligence Workbench"
TAGLINE = "Name the movement. Own the evidence."
API_VERSION = "1.0.0"
GENERATOR_VERSION = "1.0.0"
ALLOWED_SECTIONS = {
    "executive_summary",
    "kpis",
    "root_cause",
    "partners",
    "decision_log",
}
DEFAULT_SECTIONS = ["executive_summary", "kpis", "root_cause", "partners", "decision_log"]

NAVY = RGBColor(10, 24, 48)
DEEP_NAVY = RGBColor(5, 15, 31)
TEAL = RGBColor(0, 171, 169)
PALE_TEAL = RGBColor(218, 247, 245)
SKY = RGBColor(87, 180, 232)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(22, 36, 54)
MUTED = RGBColor(91, 109, 128)
LIGHT = RGBColor(242, 246, 250)
AMBER = RGBColor(238, 161, 47)
RED = RGBColor(196, 61, 74)


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _normalise_chart_axis_ids(path: Path) -> None:
    """Store chart axis IDs as unsigned integers for strict OOXML readers."""

    pattern = re.compile(rb'(<c:(?:axId|crossAx)\s+val=")-([0-9]+)(")')
    with tempfile.NamedTemporaryFile(
        suffix=".pptx", prefix="naim-axis-normalised-", dir=path.parent, delete=False
    ) as temporary:
        normalised_path = Path(temporary.name)
    try:
        with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(
            normalised_path, "w"
        ) as target:
            for member in source.infolist():
                payload = source.read(member.filename)
                if member.filename.startswith("ppt/charts/") and member.filename.endswith(".xml"):
                    payload = pattern.sub(
                        lambda match: (
                            match.group(1)
                            + str((1 << 32) - int(match.group(2))).encode("ascii")
                            + match.group(3)
                        ),
                        payload,
                    )
                target.writestr(member, payload)
        os.replace(normalised_path, path)
    finally:
        normalised_path.unlink(missing_ok=True)


def _evidence_id(run_id: str, label: str, payload: Mapping[str, Any]) -> str:
    canonical = json.dumps(
        {"run_id": run_id, "label": label, "payload": payload},
        sort_keys=True,
        separators=(",", ":"),
        default=str,
    )
    return f"EVD-{hashlib.sha256(canonical.encode()).hexdigest()[:16].upper()}"


def _format_metric(value: Any, unit: str | None) -> str:
    if value is None:
        return "N/A"
    numeric = float(value)
    if unit in {"rate", "percentage"}:
        return f"{numeric:.2%}"
    if unit in {"currency", "GBP"}:
        magnitude = abs(numeric)
        if magnitude >= 1_000_000:
            return f"£{numeric / 1_000_000:.1f}m"
        if magnitude >= 1_000:
            return f"£{numeric / 1_000:.1f}k"
        return f"£{numeric:,.0f}"
    if unit in {"bps", "basis_points"}:
        return f"{numeric:,.1f} bps"
    return f"{numeric:,.2f}"


def _set_shape_text(
    shape: Any,
    text: str,
    *,
    size: float,
    color: RGBColor,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


class DeckBuilder:
    """Small design system for an editable nAIM review deck."""

    def __init__(
        self,
        *,
        reporting_period: str,
        source_label: str,
        filter_scope: str,
        refresh_timestamp: str,
    ) -> None:
        self.presentation = Presentation()
        self.presentation.slide_width = Inches(13.333333)
        self.presentation.slide_height = Inches(7.5)
        self.reporting_period = reporting_period
        self.source_label = source_label
        self.filter_scope = filter_scope
        self.refresh_timestamp = refresh_timestamp
        self.evidence_ids: list[str] = []

    def _blank(self) -> Any:
        return self.presentation.slides.add_slide(self.presentation.slide_layouts[6])

    def _background(self, slide: Any, color: RGBColor = WHITE) -> None:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0,
            0,
            self.presentation.slide_width,
            self.presentation.slide_height,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

    def _title(self, slide: Any, title: str, subtitle: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.38), Inches(11.8), Inches(0.55))
        _set_shape_text(title_box, title, size=25, color=INK, bold=True)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.02), Inches(1.1), Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = TEAL
        accent.line.fill.background()
        if subtitle:
            box = slide.shapes.add_textbox(Inches(1.9), Inches(0.91), Inches(10.3), Inches(0.3))
            _set_shape_text(box, subtitle, size=10, color=MUTED)

    def _footer(self, slide: Any, evidence_id: str) -> None:
        self.evidence_ids.append(evidence_id)
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.04), Inches(12.05), Inches(0.015)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(207, 217, 226)
        line.line.fill.background()
        left = slide.shapes.add_textbox(Inches(0.65), Inches(7.08), Inches(8.7), Inches(0.22))
        _set_shape_text(
            left,
            f"{self.source_label}  •  Scope: {self.filter_scope}  •  Evidence: {evidence_id}",
            size=7.5,
            color=MUTED,
        )
        right = slide.shapes.add_textbox(Inches(9.35), Inches(7.08), Inches(3.35), Inches(0.22))
        _set_shape_text(
            right,
            f"Refreshed {self.refresh_timestamp[:19]}Z",
            size=7.5,
            color=MUTED,
            align=PP_ALIGN.RIGHT,
        )

    def _notes(self, slide: Any, *, talk_track: str, evidence_id: str, sources: str) -> None:
        frame = slide.notes_slide.notes_text_frame
        frame.text = (
            f"Speaker notes\n\n{talk_track}\n\n"
            f"Source block\nEvidence ID: {evidence_id}\n{sources}\n"
            f"Metric/API version: {API_VERSION}\nFilter scope: {self.filter_scope}\n"
            f"Refresh timestamp: {self.refresh_timestamp}\n"
            "Synthetic-data disclaimer: synthetic, institution-neutral demonstration data; "
            "human review is required."
        )

    def add_cover(self, *, comparison_period: str | None, run_id: str, data_mode: str) -> None:
        slide = self._blank()
        self._background(slide, DEEP_NAVY)
        mark = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(2.3), Inches(0.7))
        _set_shape_text(mark, "nAIM", size=30, color=TEAL, bold=True)
        title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.3), Inches(1.5))
        _set_shape_text(title, "Portfolio Intelligence Review", size=35, color=WHITE, bold=True)
        subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(3.55), Inches(10.6), Inches(0.55))
        comparison = f" vs {comparison_period}" if comparison_period else ""
        _set_shape_text(
            subtitle,
            f"Reporting period {self.reporting_period}{comparison}  •  {data_mode}",
            size=17,
            color=RGBColor(194, 211, 228),
        )
        tagline = slide.shapes.add_textbox(Inches(0.8), Inches(5.55), Inches(8.5), Inches(0.5))
        _set_shape_text(tagline, TAGLINE, size=18, color=TEAL, bold=True)
        disclaimer = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.7), Inches(0.35))
        _set_shape_text(
            disclaimer,
            "Synthetic, institution-neutral demonstration data • Recommendations require human review",
            size=9,
            color=RGBColor(164, 183, 204),
        )
        evidence = _evidence_id(run_id, "cover", {"period": self.reporting_period})
        self._notes(
            slide,
            talk_track=(
                "Open by naming the reporting period, active data mode, and decision purpose. "
                "Emphasise that the deck is an editable evidence pack, not an automated decision."
            ),
            evidence_id=evidence,
            sources=f"Pipeline run {run_id}; API metadata {API_VERSION}.",
        )

    def add_executive_summary(
        self,
        command: Mapping[str, Any],
        root_cause: Mapping[str, Any],
        run_id: str,
    ) -> None:
        slide = self._blank()
        self._background(slide)
        self._title(slide, "Executive summary", "Decision-ready evidence, not causal proof")
        kpis = list(command.get("kpis", []))[:4]
        for index, kpi in enumerate(kpis):
            left = 0.65 + index * 3.05
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(1.42),
                Inches(2.75),
                Inches(1.25),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT
            card.line.color.rgb = RGBColor(219, 228, 236)
            label = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(1.58), Inches(2.45), Inches(0.3)
            )
            _set_shape_text(
                label,
                str(kpi.get("name") or kpi.get("metric_name") or kpi.get("metric_id")),
                size=9,
                color=MUTED,
                bold=True,
            )
            value = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(1.95), Inches(2.45), Inches(0.48)
            )
            _set_shape_text(
                value,
                _format_metric(kpi.get("value"), kpi.get("unit")),
                size=22,
                color=INK,
                bold=True,
            )
        finding = root_cause.get("finding") or {}
        finding_text = (
            f"{finding.get('primary_driver', 'No single driver')} is the largest measured "
            f"{finding.get('primary_dimension', 'segment')} contributor to a "
            f"{float(finding.get('observed_change_bps') or 0):+.1f} bps movement. "
            "The decomposition is associational and requires investigation."
        )
        insight = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(3.05), Inches(7.9), Inches(2.9)
        )
        insight.fill.solid()
        insight.fill.fore_color.rgb = PALE_TEAL
        insight.line.color.rgb = TEAL
        header = slide.shapes.add_textbox(Inches(0.95), Inches(3.35), Inches(7.3), Inches(0.35))
        _set_shape_text(header, "Named movement", size=13, color=TEAL, bold=True)
        body = slide.shapes.add_textbox(Inches(0.95), Inches(3.85), Inches(7.2), Inches(1.65))
        _set_shape_text(body, finding_text, size=15, color=INK)
        action = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(3.05), Inches(3.85), Inches(2.9)
        )
        action.fill.solid()
        action.fill.fore_color.rgb = NAVY
        action.line.fill.background()
        action_header = slide.shapes.add_textbox(
            Inches(9.15), Inches(3.35), Inches(3.25), Inches(0.35)
        )
        _set_shape_text(action_header, "Governed next step", size=13, color=TEAL, bold=True)
        action_body = slide.shapes.add_textbox(
            Inches(9.15), Inches(3.85), Inches(3.2), Inches(1.65)
        )
        _set_shape_text(
            action_body,
            "Investigate the largest validated driver; document evidence, owner and approval before action.",
            size=14,
            color=WHITE,
        )
        evidence = _evidence_id(run_id, "executive-summary", {"kpis": kpis, "finding": finding})
        self._footer(slide, evidence)
        self._notes(
            slide,
            talk_track=(
                "Lead with the four governed KPIs, then name the measured movement. "
                "Describe attribution as associational unless the evidence contract explicitly says otherwise."
            ),
            evidence_id=evidence,
            sources="GET /api/v1/command-centre; GET /api/v1/root-cause.",
        )

    def add_kpis(self, command: Mapping[str, Any], run_id: str) -> None:
        slide = self._blank()
        self._background(slide)
        self._title(slide, "Portfolio scorecard", "Current value, prior value and movement")
        kpis = list(command.get("kpis", []))[:8]
        table_shape = slide.shapes.add_table(
            len(kpis) + 1, 4, Inches(0.65), Inches(1.4), Inches(12.05), Inches(4.95)
        )
        table = table_shape.table
        widths = [4.8, 2.2, 2.2, 2.85]
        for column, width in zip(table.columns, widths, strict=True):
            column.width = Inches(width)
        headers = ("Metric", "Current", "Prior", "Movement")
        for column_index, header in enumerate(headers):
            cell = table.cell(0, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            _set_shape_text(cell, header, size=10, color=WHITE, bold=True)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for row_index, kpi in enumerate(kpis, start=1):
            movement = None
            if kpi.get("value") is not None and kpi.get("prior_value") is not None:
                movement = float(kpi["value"]) - float(kpi["prior_value"])
            values = (
                str(kpi.get("name") or kpi.get("metric_name") or kpi.get("metric_id")),
                _format_metric(kpi.get("value"), kpi.get("unit")),
                _format_metric(kpi.get("prior_value"), kpi.get("unit")),
                _format_metric(movement, kpi.get("unit")),
            )
            for column_index, value in enumerate(values):
                cell = table.cell(row_index, column_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_index % 2 else LIGHT
                _set_shape_text(cell, value, size=9.5, color=INK, bold=column_index == 0)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        evidence = _evidence_id(run_id, "kpi-scorecard", {"kpis": kpis})
        self._footer(slide, evidence)
        self._notes(
            slide,
            talk_track=(
                "Walk down the scorecard in business order. Call out denominator and version caveats "
                "where a movement is material. Missing values remain N/A and are never backfilled from demo facts."
            ),
            evidence_id=evidence,
            sources="GET /api/v1/command-centre; governed metric registry.",
        )

    def add_root_cause(self, root_cause: Mapping[str, Any], run_id: str) -> None:
        slide = self._blank()
        self._background(slide)
        self._title(slide, "Movement bridge", "Exact additive attribution in basis points")
        finding = dict(root_cause.get("finding") or {})
        mix = float(finding.get("mix_contribution_bps") or 0.0)
        within = float(finding.get("within_segment_contribution_bps") or 0.0)
        total = float(finding.get("observed_change_bps") or mix + within)
        chart_data = ChartData()
        chart_data.categories = ["Mix", "Within segment", "Total"]
        chart_data.add_series("Movement (bps)", [mix, within, total])
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.75),
            Inches(1.45),
            Inches(7.3),
            Inches(4.9),
            chart_data,
        ).chart
        chart.has_legend = False
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.tick_labels.font.size = Pt(9)
        chart.category_axis.tick_labels.font.size = Pt(10)
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = TEAL
        explanation = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(1.45), Inches(4.35), Inches(4.9)
        )
        explanation.fill.solid()
        explanation.fill.fore_color.rgb = LIGHT
        explanation.line.color.rgb = RGBColor(219, 228, 236)
        header = slide.shapes.add_textbox(Inches(8.7), Inches(1.85), Inches(3.7), Inches(0.35))
        _set_shape_text(header, "Interpretation", size=14, color=TEAL, bold=True)
        body = slide.shapes.add_textbox(Inches(8.7), Inches(2.4), Inches(3.55), Inches(2.95))
        causal = str(finding.get("causal_status") or "ASSOCIATIONAL")
        _set_shape_text(
            body,
            (
                f"Total movement: {total:+.1f} bps\n\n"
                f"Mix effect: {mix:+.1f} bps\n"
                f"Within-segment effect: {within:+.1f} bps\n\n"
                f"Causal status: {causal}\n\n"
                "Use this bridge to target investigation; do not treat it as proof of cause."
            ),
            size=13,
            color=INK,
        )
        evidence = _evidence_id(run_id, "root-cause", finding)
        self._footer(slide, evidence)
        self._notes(
            slide,
            talk_track=(
                "Explain that mix plus within-segment effects reconcile exactly to measured movement. "
                "The bridge is descriptive attribution and the largest component should become the investigation hypothesis."
            ),
            evidence_id=evidence,
            sources="GET /api/v1/root-cause; exact decomposition module.",
        )

    def add_partners(self, partners: Mapping[str, Any], run_id: str) -> None:
        slide = self._blank()
        self._background(slide)
        self._title(
            slide, "Partner risk monitor", "Largest validated exposures and current ratings"
        )
        rows = sorted(
            list(partners.get("data", [])),
            key=lambda row: (
                float(row.get("average_balance") or 0) * float(row.get("active_accounts") or 0)
            ),
            reverse=True,
        )[:6]
        table_shape = slide.shapes.add_table(
            len(rows) + 1, 5, Inches(0.65), Inches(1.42), Inches(12.05), Inches(4.95)
        )
        table = table_shape.table
        for column, width in zip(table.columns, [2.6, 2.3, 2.2, 2.2, 2.75], strict=True):
            column.width = Inches(width)
        headers = ("Partner", "Exposure", "Loss", "Rating", "Evidence status")
        for column_index, header in enumerate(headers):
            cell = table.cell(0, column_index)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            _set_shape_text(cell, header, size=10, color=WHITE, bold=True)
        for row_index, row in enumerate(rows, start=1):
            rating = row.get("rating") or {}
            values = (
                str(row.get("partner_name") or row.get("partner_id") or "N/A"),
                _format_metric(
                    float(row.get("average_balance") or 0) * float(row.get("active_accounts") or 0),
                    "currency",
                ),
                _format_metric(
                    (
                        (
                            float(row.get("credit_loss") or 0)
                            + float(row.get("confirmed_fraud_loss") or 0)
                        )
                        / max(
                            float(row.get("average_balance") or 0)
                            * float(row.get("active_accounts") or 0),
                            1,
                        )
                    ),
                    "rate",
                ),
                str(rating.get("grade") or rating.get("rating_grade") or "N/A"),
                "Validated API row",
            )
            for column_index, value in enumerate(values):
                cell = table.cell(row_index, column_index)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if row_index % 2 else LIGHT
                _set_shape_text(cell, value, size=9.5, color=INK, bold=column_index == 0)
        evidence = _evidence_id(run_id, "partners", {"rows": rows})
        self._footer(slide, evidence)
        self._notes(
            slide,
            talk_track=(
                "Review the largest exposures first. Ratings are governed composites, not default probabilities. "
                "Any action must consider contractual terms and capacity constraints."
            ),
            evidence_id=evidence,
            sources="GET /api/v1/partners; governed rating methodology.",
        )

    def add_decision_log(self, root_cause: Mapping[str, Any], run_id: str) -> None:
        slide = self._blank()
        self._background(slide)
        self._title(slide, "Decision and control log", "Explicit ownership, evidence and approval")
        rows = [
            (
                "Investigate",
                "Validate the largest movement driver against segment and contract evidence.",
                "Portfolio Analyst",
                "Pending",
            ),
            (
                "Challenge",
                "Run a bounded scenario; record constraints and infeasibility before recommendation.",
                "Strategy Analyst",
                "Draft",
            ),
            (
                "Approve",
                "Model/rating and configuration changes require a versioned reviewer decision.",
                "Model Validator",
                "Required",
            ),
        ]
        for index, (verb, action, owner, status) in enumerate(rows):
            top = 1.45 + index * 1.55
            number = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(0.8), Inches(top + 0.1), Inches(0.65), Inches(0.65)
            )
            number.fill.solid()
            number.fill.fore_color.rgb = TEAL
            number.line.fill.background()
            _set_shape_text(
                number, str(index + 1), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER
            )
            verb_box = slide.shapes.add_textbox(
                Inches(1.75), Inches(top), Inches(1.5), Inches(0.36)
            )
            _set_shape_text(verb_box, verb, size=14, color=TEAL, bold=True)
            action_box = slide.shapes.add_textbox(
                Inches(3.1), Inches(top), Inches(6.0), Inches(0.9)
            )
            _set_shape_text(action_box, action, size=12, color=INK)
            owner_box = slide.shapes.add_textbox(
                Inches(9.35), Inches(top), Inches(1.8), Inches(0.55)
            )
            _set_shape_text(owner_box, owner, size=10, color=MUTED, bold=True)
            status_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(11.35),
                Inches(top),
                Inches(1.25),
                Inches(0.46),
            )
            status_box.fill.solid()
            status_box.fill.fore_color.rgb = (
                PALE_TEAL if status != "Required" else RGBColor(255, 241, 213)
            )
            status_box.line.fill.background()
            _set_shape_text(status_box, status, size=9, color=INK, bold=True, align=PP_ALIGN.CENTER)
        evidence = _evidence_id(
            run_id, "decision-log", {"finding": root_cause.get("finding"), "rows": rows}
        )
        self._footer(slide, evidence)
        self._notes(
            slide,
            talk_track=(
                "Close with owners and approval state. The analytical deck proposes governed next steps; "
                "it does not execute strategies, publish configurations, or approve models."
            ),
            evidence_id=evidence,
            sources="Workflow role matrix; root-cause evidence contract.",
        )

    def add_appendix(self, metadata: Mapping[str, Any], run_id: str) -> None:
        slide = self._blank()
        self._background(slide)
        self._title(slide, "Appendix — provenance and limitations", "Portable evidence block")
        items = [
            ("Pipeline run", str(metadata.get("run_id") or run_id)),
            ("Configuration hash", str(metadata.get("configuration_hash") or "N/A")),
            ("Metric/API version", API_VERSION),
            ("Reporting period", self.reporting_period),
            ("Filter scope", self.filter_scope),
            ("Refresh timestamp", self.refresh_timestamp),
            ("Data disclaimer", "Synthetic, institution-neutral demonstration data"),
            ("Decision control", "All recommendations require human review and explicit approval"),
        ]
        for index, (label, value) in enumerate(items):
            top = 1.35 + index * 0.63
            label_box = slide.shapes.add_textbox(
                Inches(0.85), Inches(top), Inches(2.35), Inches(0.35)
            )
            _set_shape_text(label_box, label, size=10, color=MUTED, bold=True)
            value_box = slide.shapes.add_textbox(
                Inches(3.2), Inches(top), Inches(9.0), Inches(0.42)
            )
            _set_shape_text(value_box, value, size=10.5, color=INK)
        evidence = _evidence_id(run_id, "appendix", {"items": items})
        self._footer(slide, evidence)
        self._notes(
            slide,
            talk_track="Use this appendix to answer provenance, versioning and limitation questions.",
            evidence_id=evidence,
            sources="GET /api/v1/metadata; artifact manifest; capability registry.",
        )


def validate_presentation(path: Path, expected_sections: Sequence[str]) -> dict[str, Any]:
    """Validate package integrity, titles, notes, object bounds and chart presence."""

    checks: list[dict[str, Any]] = []
    with zipfile.ZipFile(path) as archive:
        corrupt_member = archive.testzip()
        checks.append(
            {
                "check": "office_package_integrity",
                "status": "PASS" if corrupt_member is None else "FAIL",
                "detail": corrupt_member,
            }
        )
    presentation = Presentation(path)
    expected_count = 1 + len(expected_sections)
    checks.append(
        {
            "check": "slide_count",
            "status": "PASS" if len(presentation.slides) == expected_count else "FAIL",
            "expected": expected_count,
            "actual": len(presentation.slides),
        }
    )
    missing_titles: list[int] = []
    missing_notes: list[int] = []
    out_of_bounds: list[dict[str, Any]] = []
    empty_charts: list[int] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        visible_text = [
            shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip()
        ]
        if not visible_text:
            missing_titles.append(slide_index)
        if not slide.notes_slide.notes_text_frame.text.strip():
            missing_notes.append(slide_index)
        for shape in slide.shapes:
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > presentation.slide_width + 1
                or shape.top + shape.height > presentation.slide_height + 1
            ):
                out_of_bounds.append({"slide": slide_index, "shape": shape.name})
            if getattr(shape, "has_chart", False):
                try:
                    if not shape.chart.series:
                        empty_charts.append(slide_index)
                except (AttributeError, ValueError):
                    empty_charts.append(slide_index)
    checks.extend(
        [
            {
                "check": "missing_titles",
                "status": "PASS" if not missing_titles else "FAIL",
                "slides": missing_titles,
            },
            {
                "check": "speaker_notes",
                "status": "PASS" if not missing_notes else "FAIL",
                "slides": missing_notes,
            },
            {
                "check": "shape_bounds",
                "status": "PASS" if not out_of_bounds else "FAIL",
                "details": out_of_bounds,
            },
            {
                "check": "empty_charts",
                "status": "PASS" if not empty_charts else "FAIL",
                "slides": empty_charts,
            },
        ]
    )
    return {
        "status": "PASS" if all(check["status"] == "PASS" for check in checks) else "FAIL",
        "slide_count": len(presentation.slides),
        "checks": checks,
    }


def generate_presentation(
    service: WorkbenchService,
    payload: Mapping[str, Any],
    *,
    store: WorkflowStore,
    source_context: SourceContext,
    actor: str,
    output_root: Path | None = None,
) -> dict[str, Any]:
    """Generate, validate, manifest and persist one editable presentation job."""

    metadata = service.metadata()
    if payload.get("speaker_notes", True) is not True:
        raise ValueError("Generated review decks require full speaker notes")
    if payload.get("workspace_id"):
        service.workspace_detail(str(payload["workspace_id"]))
    if payload.get("basket_id"):
        service.basket_detail(
            str(payload["basket_id"]), requested_period=payload.get("reporting_period")
        )
    reporting_period = str(payload.get("reporting_period") or metadata.get("as_of"))
    comparison_period = payload.get("comparison_period")
    requested_sections = payload.get("selected_sections") or DEFAULT_SECTIONS
    sections = [str(section) for section in requested_sections]
    unknown_sections = sorted(set(sections) - ALLOWED_SECTIONS)
    if unknown_sections:
        raise ValueError(f"Unsupported presentation sections: {', '.join(unknown_sections)}")
    if len(sections) != len(set(sections)):
        raise ValueError("Presentation sections must not be duplicated")
    output_directory = (output_root or REPOSITORY_ROOT / "outputs" / "presentations").resolve()
    governed_root = (REPOSITORY_ROOT / "outputs").resolve()
    if output_root is None and not output_directory.is_relative_to(governed_root):
        raise ValueError("Presentation output must remain below outputs/")
    output_directory.mkdir(parents=True, exist_ok=True)
    reporting_month = reporting_period[:7].replace("-", "_")
    output_path = output_directory / f"nAIM_Portfolio_Intelligence_Review_{reporting_month}.pptx"
    manifest_path = output_path.with_suffix(".manifest.json")

    command = service.command_centre(period=reporting_period)
    root_cause = service.root_cause(period=reporting_period)
    partners = service.partners(period=reporting_period)
    filter_scope = json.dumps(
        {
            "workspace": payload.get("workspace_id"),
            "headline_scope": "all_portfolio",
            "approved_reference_basket": payload.get("basket_id"),
            "scenario": payload.get("scenario_name", "Baseline"),
        },
        sort_keys=True,
        separators=(",", ":"),
    )
    refresh_timestamp = datetime.now(UTC).isoformat()
    builder = DeckBuilder(
        reporting_period=reporting_period,
        source_label=f"{source_context.active_mode.value} • run {metadata['run_id']}",
        filter_scope=filter_scope,
        refresh_timestamp=refresh_timestamp,
    )
    builder.add_cover(
        comparison_period=str(comparison_period) if comparison_period else None,
        run_id=str(metadata["run_id"]),
        data_mode=source_context.active_mode.value,
    )
    section_methods = {
        "executive_summary": lambda: builder.add_executive_summary(
            command, root_cause, str(metadata["run_id"])
        ),
        "kpis": lambda: builder.add_kpis(command, str(metadata["run_id"])),
        "root_cause": lambda: builder.add_root_cause(root_cause, str(metadata["run_id"])),
        "partners": lambda: builder.add_partners(partners, str(metadata["run_id"])),
        "decision_log": lambda: builder.add_decision_log(root_cause, str(metadata["run_id"])),
    }
    for section in sections:
        section_methods[section]()
    if payload.get("include_appendix", True):
        (builder.add_appendix(metadata, str(metadata["run_id"])),)
        sections.append("appendix")

    with tempfile.NamedTemporaryFile(
        suffix=".pptx", prefix="naim-presentation-", dir=output_directory, delete=False
    ) as temporary:
        temporary_path = Path(temporary.name)
    try:
        builder.presentation.save(temporary_path)
        _normalise_chart_axis_ids(temporary_path)
        validation = validate_presentation(temporary_path, sections)
        if validation["status"] != "PASS":
            raise ValueError("Generated presentation failed structural validation")
        os.replace(temporary_path, output_path)
    finally:
        temporary_path.unlink(missing_ok=True)

    presentation_hash = _sha256(output_path)
    presentation_id = f"PRES-{presentation_hash[:16].upper()}"
    manifest = {
        "schema_version": "1.0.0",
        "product": PRODUCT,
        "generator_version": GENERATOR_VERSION,
        "presentation_id": presentation_id,
        "artifact": {
            "filename": output_path.name,
            "bytes": output_path.stat().st_size,
            "sha256": presentation_hash,
        },
        "built_at_utc": refresh_timestamp,
        "created_by": actor,
        "source_snapshot_id": source_context.run_id or metadata.get("run_id"),
        "data_mode": source_context.active_mode.value,
        "reporting_period": reporting_period,
        "comparison_period": comparison_period,
        "configuration_hash": source_context.configuration_hash,
        "dataset_hash": source_context.dataset_hash,
        "metric_version": API_VERSION,
        "filter_scope": json.loads(filter_scope),
        "selected_sections": sections,
        "presentation_template": payload.get("presentation_template", "executive_review"),
        "detail_level": payload.get("detail_level", "standard"),
        "commentary_length": payload.get("commentary_length", 450),
        "evidence_ids": builder.evidence_ids,
        "validation": validation,
        "synthetic_data_disclaimer": (
            "Synthetic, institution-neutral demonstration data; human review is required."
        ),
        "reconciliation": {
            "api_source": "/api/v1/command-centre",
            "excel_reference": "not_generated_in_this_request",
            "status": "API_RECONCILED",
        },
    }
    manifest_path.write_text(json.dumps(manifest, indent=2) + "\n", encoding="utf-8")
    state = {
        "record_kind": "presentation_job",
        "presentation_id": presentation_id,
        "filename": output_path.name,
        "manifest_filename": manifest_path.name,
        "status": "completed",
        "slide_count": validation["slide_count"],
        "validation_status": validation["status"],
        "reporting_period": reporting_period,
        "approval_required": True,
    }
    try:
        store.create(
            "configuration_change",
            presentation_id,
            state,
            actor=actor,
            approval_state="DRAFT",
        )
    except DuplicateObject:
        current = store.get("configuration_change", presentation_id)
        store.update(
            "configuration_change",
            presentation_id,
            state,
            expected_version=int(current["version"]),
            actor=actor,
            approval_state="DRAFT",
            replace=True,
        )
    return {
        **state,
        "download_url": f"/api/v1/presentations/{presentation_id}/download",
        "manifest_url": f"/api/v1/presentations/{presentation_id}/manifest",
        "approval_state": "DRAFT",
        "source_context": source_context.public(),
    }


def presentation_record(store: WorkflowStore, presentation_id: str) -> dict[str, Any]:
    record = store.get("configuration_change", presentation_id)
    state = dict(record["state"])
    if state.get("record_kind") != "presentation_job":
        raise KeyError(presentation_id)
    return {**state, "version": record["version"], "approval_state": record["approval_state"]}


def resolve_presentation_file(
    store: WorkflowStore,
    presentation_id: str,
    *,
    manifest: bool = False,
    output_root: Path | None = None,
) -> Path:
    state = presentation_record(store, presentation_id)
    key = "manifest_filename" if manifest else "filename"
    filename = str(state.get(key, ""))
    if not filename or Path(filename).name != filename:
        raise KeyError(presentation_id)
    root = (output_root or REPOSITORY_ROOT / "outputs" / "presentations").resolve()
    candidate = (root / filename).resolve()
    if candidate.parent != root or not candidate.is_file():
        raise KeyError(presentation_id)
    return candidate


def list_presentations(store: WorkflowStore) -> list[dict[str, Any]]:
    rows = []
    for record in store.list("configuration_change"):
        if record["state"].get("record_kind") != "presentation_job":
            continue
        rows.append(
            {
                **record["state"],
                "version": record["version"],
                "approval_state": record["approval_state"],
                "download_url": (f"/api/v1/presentations/{record['external_id']}/download"),
            }
        )
    return rows
